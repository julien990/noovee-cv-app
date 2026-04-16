import streamlit as st
import os, re, json, subprocess, shutil, time, base64, hashlib
import urllib.request
import pandas as pd
import urllib.parse
import unicodedata
from datetime import datetime, timedelta
from pptx import Presentation
from docx import Document as DocxDocument
import pypdf
import tempfile
from io import BytesIO

# ============================================================================
# 1. CONFIGURATION ET CHEMINS
# ============================================================================
NO_DATA = "---"

try:
    IS_CLOUD = "gcp_service_account" in st.secrets
except:
    IS_CLOUD = False

PATH_LOCAL = "/Users/juliensac/Library/CloudStorage/GoogleDrive-julien@miint.pro/Drive partagés/Noovee - CV"

if IS_CLOUD:
    os.makedirs("/tmp/noovee/PDF", exist_ok=True)
    PATH_DOSSIER = "/tmp/noovee"
    PATH_PDF     = "/tmp/noovee/PDF"
    PATH_DB      = "/tmp/noovee/contacts_db.json"
    PATH_BACKUP  = "/tmp/noovee/backups"
else:
    PATH_DOSSIER = PATH_LOCAL
    PATH_PDF     = os.path.join(PATH_LOCAL, "PDF")
    PATH_DB      = os.path.join(PATH_LOCAL, "contacts_db.json")
    PATH_BACKUP  = os.path.join(PATH_LOCAL, "backups")

ANTHROPIC_URL = "https://api.anthropic.com/v1/messages"
MISTRAL_URL   = "https://api.mistral.ai/v1/chat/completions"
OPENAI_URL    = "https://api.openai.com/v1/chat/completions"
GROQ_URL      = "https://api.groq.com/openai/v1/chat/completions"
GEMINI_URL    = "https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash:generateContent"

# ============================================================================
# 2. CHARGEMENT DES CLÉS IA
# ============================================================================
def charger_toutes_les_cles():
    cles = {}
    env_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), ".env")
    vars_ok = ["ANTHROPIC_API_KEY","MISTRAL_API_KEY","OPENAI_API_KEY","GROQ_API_KEY","GEMINI_API_KEY"]
    for var in vars_ok:
        val = os.environ.get(var, "")
        if val: cles[var] = val
    if os.path.exists(env_path):
        with open(env_path) as f:
            for line in f:
                line = line.strip()
                for var in vars_ok:
                    if line.startswith(var + "=") and var not in cles:
                        cles[var] = line.split("=",1)[1].strip().strip('"').strip("'")
    return cles

ALL_KEYS = charger_toutes_les_cles()

def get_provider_extraction():
    for v in ["MISTRAL_API_KEY","ANTHROPIC_API_KEY","OPENAI_API_KEY","GEMINI_API_KEY"]:
        if v in ALL_KEYS: return v.replace("_API_KEY",""), ALL_KEYS[v]
    return None, ""

def get_provider_scoring():
    for v in ["OPENAI_API_KEY","ANTHROPIC_API_KEY","MISTRAL_API_KEY","GEMINI_API_KEY"]:
        if v in ALL_KEYS: return v.replace("_API_KEY",""), ALL_KEYS[v]
    return None, ""

def appeler_ia(prompt, provider, key, max_tokens=800, model_override=None):
    if not provider or not key: raise Exception("Pas de cle IA")
    if provider == "ANTHROPIC":
        model = model_override or "claude-haiku-4-5-20251001"
        payload = json.dumps({"model":model,"max_tokens":max_tokens,
            "messages":[{"role":"user","content":prompt}]}).encode()
        req = urllib.request.Request(ANTHROPIC_URL, data=payload,
            headers={"Content-Type":"application/json","x-api-key":key,
                     "anthropic-version":"2023-06-01"}, method="POST")
        with urllib.request.urlopen(req, timeout=25) as r:
            return json.loads(r.read())["content"][0]["text"].strip()
    elif provider in ("MISTRAL","OPENAI","GROQ"):
        url   = {"MISTRAL":MISTRAL_URL,"OPENAI":OPENAI_URL,"GROQ":GROQ_URL}[provider]
        model = model_override or {"MISTRAL":"mistral-small-latest","OPENAI":"gpt-4o-mini","GROQ":"llama3-8b-8192"}[provider]
        payload = json.dumps({"model":model,"max_tokens":max_tokens,
            "messages":[{"role":"user","content":prompt}],"temperature":0}).encode()
        req = urllib.request.Request(url, data=payload,
            headers={"Content-Type":"application/json","Authorization":"Bearer "+key}, method="POST")
        with urllib.request.urlopen(req, timeout=25) as r:
            return json.loads(r.read())["choices"][0]["message"]["content"].strip()
    elif provider == "GEMINI":
        payload = json.dumps({"contents":[{"parts":[{"text":prompt}]}]}).encode()
        req = urllib.request.Request(GEMINI_URL+"?key="+key, data=payload,
            headers={"Content-Type":"application/json"}, method="POST")
        with urllib.request.urlopen(req, timeout=25) as r:
            return json.loads(r.read())["candidates"][0]["content"]["parts"][0]["text"].strip()
    raise Exception("Provider inconnu: "+str(provider))

_first = list(ALL_KEYS.items())[0] if ALL_KEYS else (None, "")
AI_PROVIDER = _first[0]
AI_KEY      = _first[1]

# ============================================================================
# 3. NORMALISATION
# ============================================================================
def normaliser(texte):
    if not texte:
        return ""
    t = unicodedata.normalize("NFD", texte.lower())
    t = "".join(c for c in t if unicodedata.category(c) != "Mn")
    for c in ["\u2019", "\u2018", "\u2032", "'", "`"]:
        t = t.replace(c, " ")
    t = t.replace("-", " ").replace("_", " ")
    t = re.sub(r"[^\w\s]", " ", t)
    return re.sub(r"\s+", " ", t).strip()

# ============================================================================
# 4. EXTRACTION DE TEXTE DES FICHIERS
# ============================================================================
def extraire_texte_pptx(filepath):
    try:
        prs = Presentation(filepath)
        texte = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texte.append(shape.text)
        return "\n".join(texte)
    except Exception as e:
        return f"Erreur extraction PPTX: {str(e)}"

def extraire_texte_docx(filepath):
    try:
        doc = DocxDocument(filepath)
        texte = []
        for para in doc.paragraphs:
            if para.text.strip():
                texte.append(para.text)
        return "\n".join(texte)
    except Exception as e:
        return f"Erreur extraction DOCX: {str(e)}"

def extraire_texte_pdf(filepath):
    try:
        texte = []
        with open(filepath, 'rb') as f:
            reader = pypdf.PdfReader(f)
            for page in reader.pages:
                text = page.extract_text()
                if text.strip():
                    texte.append(text)
        return "\n".join(texte)
    except Exception as e:
        return f"Erreur extraction PDF: {str(e)}"

def extraire_texte(filepath, extension):
    if extension == ".pptx":
        return extraire_texte_pptx(filepath)
    elif extension == ".docx":
        return extraire_texte_docx(filepath)
    elif extension == ".pdf":
        return extraire_texte_pdf(filepath)
    else:
        return ""

# ============================================================================
# 5. CONVERSION EN PDF
# ============================================================================
def convertir_en_pdf(input_path, output_path):
    try:
        soffice = None
        for c in ["/Applications/LibreOffice.app/Contents/MacOS/soffice",
                  shutil.which("soffice"), shutil.which("libreoffice")]:
            if c and os.path.exists(c):
                soffice = c
                break
        
        if not soffice:
            return False, "LibreOffice non installé"
        
        subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf", "--outdir", os.path.dirname(output_path), input_path],
            capture_output=True, timeout=30
        )
        return True, "PDF généré"
    except Exception as e:
        return False, str(e)

# ============================================================================
# 6. EXTRACTION D'INFOS PAR IA (COMPLÈTE)
# ============================================================================
def extraire_infos_cv_complete(texte_cv):
    """Extraction complète avec Mistral + Claude"""
    
    prompt_extraction = f"""Vous êtes un expert RH avec 20 ans d'expérience. Analysez ce CV en détail et extrayez les informations structurées suivantes en JSON valide.

INSTRUCTIONS STRICTES:
- Cherchez TOUTES les informations disponibles
- Si une info n'existe pas, utilisez "---"
- Les compétences: lista complète (10-20 items)
- Les formations: format "Diplôme - École - Année"
- Les certifications: format "Certification - Organisme - Année"
- Salaire attendu: format "XX000€" ou "---"
- Disponibilité: "Immédiate", "2 semaines", "1 mois", "À négocier" ou "---"

Répondez UNIQUEMENT en JSON valide, sans texte additionnel:

{{
    "nom": "...",
    "prenom": "...",
    "email": "...",
    "telephone": "...",
    "linkedin": "...",
    "github": "...",
    "secteur": "...",
    "poste": "...",
    "competences": ["...", "..."],
    "experience_ans": "...",
    "location": "...",
    "formations": ["...", "..."],
    "certifications": ["...", "..."],
    "salaire_attendu": "...",
    "disponibilite": "...",
    "langues": ["...", "..."],
    "notes": "..."
}}

CV à analyser:
{texte_cv[:4000]}
"""

    try:
        mistral_provider, mistral_key = get_provider_extraction()
        if mistral_provider and mistral_provider == "MISTRAL":
            resultat_mistral = appeler_ia(prompt_extraction, mistral_provider, mistral_key, max_tokens=1000)
            try:
                return json.loads(resultat_mistral)
            except:
                pass
        
        claude_provider = "ANTHROPIC"
        claude_key = ALL_KEYS.get("ANTHROPIC_API_KEY")
        if claude_key:
            resultat_claude = appeler_ia(prompt_extraction, claude_provider, claude_key, max_tokens=1000)
            try:
                return json.loads(resultat_claude)
            except:
                pass
        
        return {
            "nom": "---", "prenom": "---", "email": "---", "telephone": "---",
            "linkedin": "---", "github": "---", "secteur": "---", "poste": "---",
            "competences": [], "experience_ans": "---", "location": "---",
            "formations": [], "certifications": [], "salaire_attendu": "---",
            "disponibilite": "---", "langues": [], "notes": "Extraction IA échouée"
        }
    except Exception as e:
        st.error(f"Erreur IA: {str(e)}")
        return {}

# ============================================================================
# 7. BASE DE DONNÉES
# ============================================================================
def charger_db():
    if os.path.exists(PATH_DB):
        try:
            with open(PATH_DB, 'r', encoding='utf-8') as f:
                return json.load(f)
        except:
            return {}
    return {}

def sauvegarder_db(db):
    os.makedirs(os.path.dirname(PATH_DB), exist_ok=True)
    
    # Créer un backup
    os.makedirs(PATH_BACKUP, exist_ok=True)
    backup_name = f"backup_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json"
    backup_path = os.path.join(PATH_BACKUP, backup_name)
    if os.path.exists(PATH_DB):
        shutil.copy(PATH_DB, backup_path)
    
    # Sauvegarder la DB
    with open(PATH_DB, 'w', encoding='utf-8') as f:
        json.dump(db, f, ensure_ascii=False, indent=2)

def creer_fiche_contact(base_name, infos):
    """Crée une fiche contact complète avec métadonnées"""
    return {
        **infos,
        "id": base_name,
        "date_ajout": datetime.now().isoformat(),
        "date_modif": datetime.now().isoformat(),
        "tags": [],
        "favoris": False,
        "score_global": 0,
        "ia_enrichi": True,
        "whatsapp_envoye": False,
        "notes_personnelles": "",
        "historique": [],
        "matches": []
    }

# ============================================================================
# 8. SCORING ET MATCHING
# ============================================================================
def scorer_cv(competences_cv, competences_offre):
    """Scoring simple basé sur la correspondance de compétences"""
    if not competences_cv or not competences_offre:
        return 0
    
    correspondances = sum(1 for comp in competences_cv if any(normaliser(comp) in normaliser(off) or normaliser(off) in normaliser(comp) for off in competences_offre))
    score = min(100, int((correspondances / len(competences_offre)) * 100))
    return score

# ============================================================================
# 9. EXPORTS
# ============================================================================
def exporter_csv(db):
    """Exporte la DB en CSV"""
    rows = []
    for base_name, info in db.items():
        rows.append({
            "ID": base_name,
            "Nom": info.get("nom", ""),
            "Prénom": info.get("prenom", ""),
            "Email": info.get("email", ""),
            "Téléphone": info.get("telephone", ""),
            "Secteur": info.get("secteur", ""),
            "Poste": info.get("poste", ""),
            "Expérience": info.get("experience_ans", ""),
            "Localisation": info.get("location", ""),
            "Compétences": ", ".join(info.get("competences", [])),
            "LinkedIn": info.get("linkedin", ""),
            "Disponibilité": info.get("disponibilite", ""),
            "Date ajout": info.get("date_ajout", "")
        })
    
    return pd.DataFrame(rows)

# ============================================================================
# 10. CONFIGURATION STREAMLIT
# ============================================================================
st.set_page_config(
    page_title="Noovee - CV Manager Pro",
    page_icon="📄",
    layout="wide",
    initial_sidebar_state="expanded"
)

st.markdown("""
<style>
    .main { padding: 0 2rem; }
    .stTabs [data-baseweb="tab-list"] { gap: 0.5rem; }
    .stButton>button { border-radius: 0.5rem; font-weight: 600; }
    .cv-card { 
        border: 2px solid #e0e0e0; 
        border-left: 4px solid #2196F3;
        border-radius: 0.5rem; 
        padding: 1.5rem; 
        margin: 1rem 0;
        background: linear-gradient(135deg, #f5f7fa 0%, #c3cfe2 100%);
        transition: all 0.3s ease;
    }
    .cv-card:hover {
        border-left-color: #FF6B6B;
        box-shadow: 0 4px 12px rgba(0,0,0,0.1);
    }
    .metric-card {
        background: #f0f2f6;
        border-radius: 0.5rem;
        padding: 1rem;
        text-align: center;
    }
    .success-badge { background: #d4edda; color: #155724; padding: 0.5rem 1rem; border-radius: 0.25rem; font-weight: bold; }
    .warning-badge { background: #fff3cd; color: #856404; padding: 0.5rem 1rem; border-radius: 0.25rem; font-weight: bold; }
</style>
""", unsafe_allow_html=True)

# ============================================================================
# 11. INTERFACE PRINCIPAL
# ============================================================================
st.title("📄 Noovee - CV Manager Pro")
st.subheader("Plateforme intelligente de gestion de CV et matching automatique")

db = charger_db()

# Onglets
tab_ajouter, tab_cvs, tab_recherche, tab_scoring, tab_outils, tab_export = st.tabs([
    "➕ Ajouter CV", "📋 Mes CVs", "🔍 Recherche", "⭐ Scoring", "🛠 Outils", "📊 Export"
])

# ============================================================================
# ONGLET 1 — AJOUTER UN CV
# ============================================================================
with tab_ajouter:
    st.subheader("Ajouter un nouveau CV")
    
    col1, col2, col3 = st.columns([2, 1, 1])
    
    with col1:
        st.info("📤 Uploader un CV au format PPTX, DOCX ou PDF")
        uploaded_file = st.file_uploader("Sélectionnez un fichier", type=["pptx", "docx", "pdf"], key="upload_cv")
    
    with col2:
        st.metric("CVs en BD", len(db))
    
    with col3:
        st.metric("Backups", len([f for f in os.listdir(PATH_BACKUP) if f.startswith("backup_")] if os.path.exists(PATH_BACKUP) else []))
    
    if uploaded_file is not None:
        with st.spinner("⏳ Traitement du fichier..."):
            filename = uploaded_file.name
            file_ext = os.path.splitext(filename)[1].lower()
            
            temp_dir = tempfile.mkdtemp()
            temp_path = os.path.join(temp_dir, filename)
            with open(temp_path, 'wb') as f:
                f.write(uploaded_file.read())
            
            # Extraction texte
            st.info("1️⃣ Extraction du texte...")
            texte_cv = extraire_texte(temp_path, file_ext)
            
            if texte_cv and "Erreur" not in texte_cv:
                # Extraction IA
                st.info("2️⃣ Analyse par IA...")
                infos_extraites = extraire_infos_cv_complete(texte_cv)
                
                if infos_extraites:
                    base_name = re.sub(r'[^a-zA-Z0-9_]', '', filename.split('.')[0])
                    if not base_name:
                        base_name = f"cv_{int(time.time())}"
                    
                    st.success("✅ Infos extraites ! Validez les données.")
                    st.divider()
                    
                    with st.form("valider_cv", clear_on_submit=True):
                        col_f1, col_f2, col_f3 = st.columns(3)
                        
                        with col_f1:
                            nom = st.text_input("Nom", value=infos_extraites.get("nom", ""))
                            prenom = st.text_input("Prénom", value=infos_extraites.get("prenom", ""))
                            email = st.text_input("Email", value=infos_extraites.get("email", ""))
                            telephone = st.text_input("Téléphone", value=infos_extraites.get("telephone", ""))
                        
                        with col_f2:
                            secteur = st.text_input("Secteur", value=infos_extraites.get("secteur", ""))
                            poste = st.text_input("Poste", value=infos_extraites.get("poste", ""))
                            experience_ans = st.text_input("Années d'expérience", value=infos_extraites.get("experience_ans", ""))
                            location = st.text_input("Localisation", value=infos_extraites.get("location", ""))
                        
                        with col_f3:
                            linkedin = st.text_input("LinkedIn", value=infos_extraites.get("linkedin", ""))
                            github = st.text_input("GitHub", value=infos_extraites.get("github", ""))
                            salaire_attendu = st.text_input("Salaire attendu", value=infos_extraites.get("salaire_attendu", ""))
                            disponibilite = st.selectbox("Disponibilité", ["Immédiate", "2 semaines", "1 mois", "À négocier", "---"], index=4)
                        
                        competences_text = st.text_area("Compétences (séparées par des virgules)", 
                            value=", ".join(infos_extraites.get("competences", [])), height=60)
                        
                        formations_text = st.text_area("Formations (une par ligne)", 
                            value="\n".join(infos_extraites.get("formations", [])), height=60)
                        
                        certifications_text = st.text_area("Certifications (une par ligne)", 
                            value="\n".join(infos_extraites.get("certifications", [])), height=60)
                        
                        langues_text = st.text_area("Langues (séparées par des virgules)", 
                            value=", ".join(infos_extraites.get("langues", [])), height=40)
                        
                        notes = st.text_area("Notes supplémentaires", 
                            value=infos_extraites.get("notes", ""), height=80)
                        
                        col_btn1, col_btn2, col_btn3 = st.columns([1, 1, 2])
                        
                        with col_btn1:
                            submit = st.form_submit_button("✅ Valider", use_container_width=True)
                        
                        with col_btn2:
                            annuler = st.form_submit_button("❌ Annuler", use_container_width=True)
                        
                        if submit and nom and email:
                            # Préparer les données
                            competences_list = [c.strip() for c in competences_text.split(",") if c.strip()]
                            formations_list = [f.strip() for f in formations_text.split("\n") if f.strip()]
                            certifications_list = [c.strip() for c in certifications_text.split("\n") if c.strip()]
                            langues_list = [l.strip() for l in langues_text.split(",") if l.strip()]
                            
                            infos_finales = {
                                "nom": nom,
                                "prenom": prenom,
                                "email": email,
                                "telephone": telephone,
                                "linkedin": linkedin,
                                "github": github,
                                "secteur": secteur,
                                "poste": poste,
                                "competences": competences_list,
                                "experience_ans": experience_ans,
                                "location": location,
                                "formations": formations_list,
                                "certifications": certifications_list,
                                "salaire_attendu": salaire_attendu,
                                "disponibilite": disponibilite,
                                "langues": langues_list,
                                "notes": notes
                            }
                            
                            # Créer la fiche
                            nouvelle_entree = creer_fiche_contact(base_name, infos_finales)
                            
                            # Sauvegarder
                            db[base_name] = nouvelle_entree
                            sauvegarder_db(db)
                            
                            # Copier fichiers
                            dest_path = os.path.join(PATH_DOSSIER, f"{base_name}{file_ext}")
                            os.makedirs(PATH_DOSSIER, exist_ok=True)
                            shutil.copy(temp_path, dest_path)
                            
                            if file_ext != ".pdf":
                                pdf_path = os.path.join(PATH_PDF, f"{base_name}.pdf")
                                os.makedirs(PATH_PDF, exist_ok=True)
                                success, msg = convertir_en_pdf(dest_path, pdf_path)
                            
                            st.success(f"✅ CV '{prenom} {nom}' ajouté avec succès !")
                            st.balloons()
                            time.sleep(2)
                            st.rerun()
                        
                        if annuler:
                            st.info("Annulation")
            
            shutil.rmtree(temp_dir, ignore_errors=True)

# ============================================================================
# ONGLET 2 — MES CVS
# ============================================================================
with tab_cvs:
    st.subheader("📋 Liste des CVs indexés")
    
    col_search, col_count = st.columns([3, 1])
    with col_search:
        search_term = st.text_input("🔍 Rechercher...", placeholder="Nom, email, compétence, secteur...")
    with col_count:
        st.metric("Total", len(db))
    
    if db:
        # Trier par date de modification
        sorted_db = sorted(db.items(), key=lambda x: x[1].get("date_modif", ""), reverse=True)
        
        for base_name, info in sorted_db:
            full_name = f"{info.get('prenom', '')} {info.get('nom', '')}".strip()
            search_text = f"{full_name} {info.get('email', '')} {info.get('secteur', '')} {' '.join(info.get('competences', []))}".lower()
            
            if search_term.lower() in search_text:
                with st.container():
                    col_info, col_score, col_actions = st.columns([3, 1, 1])
                    
                    with col_info:
                        st.markdown(f"**{full_name}**")
                        st.caption(f"📧 {info.get('email', 'N/A')} | 💼 {info.get('poste', 'N/A')} | 🏢 {info.get('secteur', 'N/A')} | 📍 {info.get('location', 'N/A')}")
                        competences_display = ", ".join(info.get('competences', [])[:5])
                        st.caption(f"🎯 {competences_display}..." if len(info.get('competences', [])) > 5 else f"🎯 {competences_display}")
                    
                    with col_score:
                        st.metric("Score", f"{info.get('score_global', 0)}%")
                    
                    with col_actions:
                        if st.button("📋", key=f"view_{base_name}"):
                            st.session_state[f"view_{base_name}"] = True
                    
                    # Détails (expandable)
                    if st.session_state.get(f"view_{base_name}"):
                        with st.expander("Détails complets", expanded=True):
                            col_det1, col_det2 = st.columns(2)
                            with col_det1:
                                st.write(f"**Formations:** {', '.join(info.get('formations', ['Aucune']))}")
                                st.write(f"**Certifications:** {', '.join(info.get('certifications', ['Aucune']))}")
                                st.write(f"**Langues:** {', '.join(info.get('langues', ['N/A']))}")
                            with col_det2:
                                st.write(f"**Disponibilité:** {info.get('disponibilite', '---')}")
                                st.write(f"**Salaire attendu:** {info.get('salaire_attendu', '---')}")
                                st.write(f"**LinkedIn:** {info.get('linkedin', '---')}")
                            
                            st.write(f"**Notes:** {info.get('notes', 'N/A')}")
                            
                            col_btn1, col_btn2, col_btn3 = st.columns(3)
                            with col_btn1:
                                if st.button("⭐ Favori", key=f"fav_{base_name}"):
                                    db[base_name]["favoris"] = not db[base_name]["favoris"]
                                    sauvegarder_db(db)
                                    st.rerun()
                            with col_btn2:
                                if st.button("🗑 Supprimer", key=f"del_{base_name}"):
                                    del db[base_name]
                                    sauvegarder_db(db)
                                    st.rerun()
                            with col_btn3:
                                st.write(f"Ajouté: {info.get('date_ajout', 'N/A')[:10]}")
                    
                    st.divider()
    else:
        st.info("Aucun CV. Commencez par ajouter un CV ! ➕")

# ============================================================================
# ONGLET 3 — RECHERCHE
# ============================================================================
with tab_recherche:
    st.subheader("🔍 Recherche avancée")
    
    col_comp, col_sect = st.columns(2)
    
    with col_comp:
        competence = st.text_input("Compétence", placeholder="Ex: Python, Management, RGPD...")
    
    with col_sect:
        secteur_filter = st.text_input("Secteur", placeholder="Ex: Tech, Finance, Retail...")
    
    if competence or secteur_filter:
        results = []
        for base_name, info in db.items():
            comp_match = competence.lower() in " ".join(info.get("competences", [])).lower() if competence else True
            sect_match = secteur_filter.lower() in info.get("secteur", "").lower() if secteur_filter else True
            
            if comp_match and sect_match:
                results.append((base_name, info))
        
        if results:
            st.success(f"✅ {len(results)} CV(s) trouvé(s)")
            for base_name, r in results:
                st.write(f"**{r.get('prenom')} {r.get('nom')}** - {r.get('poste')} ({r.get('secteur')}) | 📧 {r.get('email')}")
        else:
            st.info("Aucun résultat")

# ============================================================================
# ONGLET 4 — SCORING
# ============================================================================
with tab_scoring:
    st.subheader("⭐ Scoring & Matching")
    
    st.write("Définissez un appel d'offres et matchez avec les CVs")
    
    offre_competences = st.text_area("Compétences requises (séparées par des virgules)", 
        placeholder="Python, Management, RGPD, Agile...")
    
    if offre_competences and db:
        comp_list = [c.strip() for c in offre_competences.split(",")]
        
        scores = []
        for base_name, info in db.items():
            score = scorer_cv(info.get("competences", []), comp_list)
            scores.append({
                "Candidat": f"{info.get('prenom')} {info.get('nom')}",
                "Score": f"{score}%",
                "Poste": info.get("poste", ""),
                "Email": info.get("email", "")
            })
        
        scores_df = pd.DataFrame(scores).sort_values("Score", ascending=False, key=lambda x: x.str.rstrip('%').astype(int))
        st.dataframe(scores_df, use_container_width=True, hide_index=True)

# ============================================================================
# ONGLET 5 — OUTILS
# ============================================================================
with tab_outils:
    st.subheader("🛠 Maintenance et configuration")
    
    col1, col2, col3, col4 = st.columns(4)
    col1.metric("CVs indexés", len(db))
    col2.metric("Dernière maj", datetime.now().strftime("%d/%m"))
    col3.metric("Backups", len([f for f in os.listdir(PATH_BACKUP) if f.startswith("backup_")] if os.path.exists(PATH_BACKUP) else []))
    
    st.divider()
    
    col_btn1, col_btn2, col_btn3 = st.columns(3)
    
    with col_btn1:
        if st.button("🔄 Actualiser", use_container_width=True):
            st.rerun()
    
    with col_btn2:
        if st.button("💾 Forcer sauvegarde", use_container_width=True):
            sauvegarder_db(db)
            st.success("✅ Sauvegardé")
    
    with col_btn3:
        if st.button("📋 Nettoyer les backups", use_container_width=True):
            if os.path.exists(PATH_BACKUP):
                backups = sorted([f for f in os.listdir(PATH_BACKUP) if f.startswith("backup_")])
                if len(backups) > 10:
                    for backup in backups[:-10]:
                        os.remove(os.path.join(PATH_BACKUP, backup))
                    st.success(f"✅ {len(backups) - 10} anciens backups supprimés")
    
    st.divider()
    
    # IA Config
    if ALL_KEYS:
        st.subheader("🤖 Fournisseurs IA")
        labels = {
            "MISTRAL_API_KEY": "🇫🇷 Mistral",
            "ANTHROPIC_API_KEY": "🤖 Claude (Anthropic)",
            "OPENAI_API_KEY": "⚡ GPT-4o (OpenAI)",
            "GROQ_API_KEY": "Groq",
            "GEMINI_API_KEY": "Gemini",
        }
        for var, key in ALL_KEYS.items():
            st.success(labels.get(var, var) + " — ✅ Connecté")
    else:
        st.warning("❌ Aucune IA configurée. Ajoutez des clés dans .env")

# ============================================================================
# ONGLET 6 — EXPORT
# ============================================================================
with tab_export:
    st.subheader("📊 Exports")
    
    if db:
        # Exporter CSV
        st.write("**Exporter en CSV**")
        df_export = exporter_csv(db)
        
        csv_buffer = BytesIO()
        df_export.to_csv(csv_buffer, index=False, encoding='utf-8')
        csv_bytes = csv_buffer.getvalue()
        
        st.download_button(
            label="📥 Télécharger CSV",
            data=csv_bytes,
            file_name=f"noovee_export_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv",
            mime="text/csv",
            use_container_width=True
        )
        
        st.divider()
        
        # Exporter JSON
        st.write("**Exporter en JSON**")
        json_str = json.dumps(db, ensure_ascii=False, indent=2)
        
        st.download_button(
            label="📥 Télécharger JSON",
            data=json_str,
            file_name=f"noovee_export_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json",
            mime="application/json",
            use_container_width=True
        )
        
        st.divider()
        
        # Aperçu
        st.write("**Aperçu des données**")
        st.dataframe(df_export, use_container_width=True)
    else:
        st.info("Aucun CV à exporter")

# ============================================================================
# PIED DE PAGE
# ============================================================================
st.divider()
col_footer1, col_footer2, col_footer3 = st.columns(3)
with col_footer1:
    st.caption(f"📁 DB: {PATH_DB}")
with col_footer2:
    st.caption(f"🏗 Path: {PATH_DOSSIER}")
with col_footer3:
    st.caption(f"📦 Noovee v2.0")
