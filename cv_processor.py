# cv_processor.py

import fitz
import subprocess
import shutil
import os
from pathlib import Path
from typing import Optional

from config import EXTRACTION_SYSTEM_PROMPT, EXTRACTION_USER_TEMPLATE, MAX_CV_TEXT_CHARS, CV_STORAGE_PATH
from ai_providers import call_ai_json

SUPABASE_URL = os.getenv("SUPABASE_URL")
SUPABASE_KEY = os.getenv("SUPABASE_KEY")
BUCKET       = "cvs"


# ── Supabase Storage ───────────────────────────────────────────────────────────

def get_supabase():
    if not SUPABASE_URL or not SUPABASE_KEY:
        return None
    try:
        from supabase import create_client
        return create_client(SUPABASE_URL, SUPABASE_KEY)
    except Exception:
        return None


def upload_to_supabase(file_bytes: bytes, filename: str) -> bool:
    """Upload un fichier vers Supabase Storage. Retourne True si succes."""
    client = get_supabase()
    if not client:
        return False
    try:
        # Supprime l'ancien si existant
        client.storage.from_(BUCKET).remove([filename])
    except Exception:
        pass
    try:
        client.storage.from_(BUCKET).upload(
            path=filename,
            file=file_bytes,
            file_options={"content-type": "application/pdf", "upsert": "true"},
        )
        return True
    except Exception as e:
        print(f"Supabase upload error: {e}")
        return False


def get_supabase_url(filename: str) -> Optional[str]:
    """Retourne l'URL publique du fichier dans Supabase Storage."""
    if not SUPABASE_URL:
        return None
    return f"{SUPABASE_URL}/storage/v1/object/public/{BUCKET}/{filename}"


def delete_from_supabase(filename: str):
    client = get_supabase()
    if not client: return
    try:
        client.storage.from_(BUCKET).remove([filename])
    except Exception:
        pass


# ── Stockage local (fallback) ──────────────────────────────────────────────────

def ensure_storage() -> Path:
    path = Path(CV_STORAGE_PATH)
    path.mkdir(parents=True, exist_ok=True)
    return path


def save_cv_file(uploaded_file) -> str:
    """Sauvegarde le fichier localement et l'uploade sur Supabase si disponible."""
    storage = ensure_storage()
    dest    = storage / uploaded_file.name
    counter = 1
    stem, suffix = Path(uploaded_file.name).stem, Path(uploaded_file.name).suffix
    while dest.exists():
        dest = storage / f"{stem}_{counter}{suffix}"
        counter += 1

    file_bytes = bytes(uploaded_file.getbuffer())

    # Sauvegarde locale
    with open(dest, "wb") as f:
        f.write(file_bytes)

    # Upload Supabase (PDF uniquement)
    if suffix.lower() == ".pdf":
        upload_to_supabase(file_bytes, dest.name)

    return dest.name


def delete_cv_file(filename: str):
    """Supprime le fichier localement et sur Supabase."""
    path = Path(CV_STORAGE_PATH) / filename
    if path.exists():
        path.unlink()
    delete_from_supabase(filename)


# ── Conversion LibreOffice ─────────────────────────────────────────────────────

LIBREOFFICE_PATHS = [
    "/Applications/LibreOffice.app/Contents/MacOS/soffice",
    "soffice",
    "libreoffice",
]

def get_libreoffice_bin() -> Optional[str]:
    for path in LIBREOFFICE_PATHS:
        if shutil.which(path) or Path(path).exists():
            return path
    return None


def convert_to_pdf(input_path: Path) -> Optional[Path]:
    soffice = get_libreoffice_bin()
    if not soffice:
        return None
    output_dir = input_path.parent
    try:
        subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(output_dir), str(input_path)],
            capture_output=True, text=True, timeout=30
        )
        pdf_path = output_dir / (input_path.stem + ".pdf")
        return pdf_path if pdf_path.exists() else None
    except Exception:
        return None


# ── Extraction texte ───────────────────────────────────────────────────────────

def extract_text_from_pdf(file_bytes: bytes) -> str:
    parts = []
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    for page in doc:
        parts.append(page.get_text("text"))
    doc.close()
    return "\n".join(parts).strip()[:MAX_CV_TEXT_CHARS]


def extract_text_from_file(file_bytes: bytes, filename: str) -> str:
    ext = Path(filename).suffix.lower()

    if ext == ".pdf":
        return extract_text_from_pdf(file_bytes)

    elif ext in (".docx", ".doc", ".pptx", ".ppt"):
        tmp_dir  = Path("/tmp/noovee_convert")
        tmp_dir.mkdir(exist_ok=True)
        tmp_file = tmp_dir / filename

        with open(tmp_file, "wb") as f:
            f.write(file_bytes)

        pdf_path = convert_to_pdf(tmp_file)

        if pdf_path and pdf_path.exists():
            with open(pdf_path, "rb") as f:
                pdf_bytes = f.read()
            text = extract_text_from_pdf(pdf_bytes)
            # Upload le PDF converti sur Supabase
            pdf_name = tmp_file.stem + ".pdf"
            upload_to_supabase(pdf_bytes, pdf_name)
            tmp_file.unlink(missing_ok=True)
            pdf_path.unlink(missing_ok=True)
            return text
        else:
            tmp_file.unlink(missing_ok=True)
            return extract_text_fallback(file_bytes, ext)
    else:
        raise ValueError(f"Format non supporte : {ext}")


def extract_text_fallback(file_bytes: bytes, ext: str) -> str:
    import io
    if ext in (".docx", ".doc"):
        from docx import Document
        doc   = Document(io.BytesIO(file_bytes))
        parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip(): parts.append(cell.text.strip())
        return "\n".join(parts)[:MAX_CV_TEXT_CHARS]
    elif ext in (".pptx", ".ppt"):
        from pptx import Presentation
        prs   = Presentation(io.BytesIO(file_bytes))
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
        return "\n".join(parts)[:MAX_CV_TEXT_CHARS]
    return ""


# ── Analyse IA ─────────────────────────────────────────────────────────────────

def analyze_cv_with_ai(cv_text: str) -> tuple:
    if not cv_text.strip():
        raise ValueError("Fichier vide ou non lisible.")
    user_prompt = EXTRACTION_USER_TEMPLATE.format(cv_text=cv_text)
    data, provider = call_ai_json(EXTRACTION_SYSTEM_PROMPT, user_prompt, preferred_provider="mistral")
    return _clean(data), provider


def _clean(data: dict) -> dict:
    from config import MAX_COMPETENCES, MAX_DOMAINES, DOMAINES

    for field in ("nom", "prenom", "email", "telephone", "poste"):
        val = data.get(field)
        data[field] = val.strip() if isinstance(val, str) and val.strip() else None

    try:
        data["annees_experience"] = max(0, int(data.get("annees_experience") or 0))
    except Exception:
        data["annees_experience"] = 0

    comp = data.get("competences", [])
    data["competences"] = [str(c).strip() for c in (comp if isinstance(comp, list) else [])][:MAX_COMPETENCES]

    dom = data.get("domaines_fonctionnels", [])
    data["domaines_fonctionnels"] = [d for d in (dom if isinstance(dom, list) else []) if d in DOMAINES][:MAX_DOMAINES]

    ent = data.get("entreprises", [])
    data["entreprises"] = [
        {"nom": str(e.get("nom","")).strip(), "secteur": str(e.get("secteur","")).strip(),
         "annees": max(0, int(e.get("annees",0) or 0))}
        for e in (ent if isinstance(ent, list) else []) if isinstance(e, dict)
    ]

    exp = data.get("experiences", [])
    data["experiences"] = [
        {"poste": str(x.get("poste","")).strip(), "entreprise": str(x.get("entreprise","")).strip(),
         "domaine": str(x.get("domaine","")).strip(), "annees": max(0, int(x.get("annees",0) or 0)),
         "mots_cles": [str(m).strip() for m in (x.get("mots_cles",[]) if isinstance(x.get("mots_cles"), list) else [])]}
        for x in (exp if isinstance(exp, list) else []) if isinstance(x, dict)
    ]
    return data


# ── Pipeline upload ────────────────────────────────────────────────────────────

def process_uploaded_cv(uploaded_file) -> tuple:
    filename   = save_cv_file(uploaded_file)
    file_bytes = bytes(uploaded_file.getbuffer())
    cv_text    = extract_text_from_file(file_bytes, uploaded_file.name)

    if not cv_text.strip():
        raise ValueError("Le fichier ne contient pas de texte lisible.")

    data, provider = analyze_cv_with_ai(cv_text)
    data["texte_brut"]  = cv_text
    data["cv_filename"] = filename
    return data, filename, provider


# ── Affichage PDF ──────────────────────────────────────────────────────────────

def get_pdf_url(filename: str) -> Optional[str]:
    """
    Retourne l'URL pour afficher le PDF :
    1. Supabase Storage (prioritaire, persiste)
    2. Fichier local (fallback)
    """
    # Essaie Supabase d'abord
    url = get_supabase_url(filename)
    if url:
        return url
    # Fallback local
    path = Path(CV_STORAGE_PATH) / filename
    if path.exists():
        return None  # Sera gere par le lecteur local
    return None


# ── Scan automatique ───────────────────────────────────────────────────────────

def scan_and_import_new_cvs() -> dict:
    import database as db
    storage   = ensure_storage()
    tracked   = db.get_tracked_filenames()
    all_files = [
        f for f in storage.iterdir()
        if f.suffix.lower() in (".pdf", ".docx", ".doc", ".pptx", ".ppt")
        and f.name not in tracked
    ]
    report = {"imported": [], "failed": [], "skipped": len(tracked)}

    for file_path in sorted(all_files):
        try:
            with open(file_path, "rb") as f:
                file_bytes = f.read()
            cv_text = extract_text_from_file(file_bytes, file_path.name)
            if not cv_text.strip():
                report["failed"].append({"filename": file_path.name, "error": "Fichier vide"})
                continue
            data, provider = analyze_cv_with_ai(cv_text)
            data["texte_brut"]  = cv_text
            data["cv_filename"] = file_path.name
            db.insert_contact(data)
            # Upload sur Supabase si PDF
            if file_path.suffix.lower() == ".pdf":
                upload_to_supabase(file_bytes, file_path.name)
            name = f"{data.get('prenom') or ''} {data.get('nom') or ''}".strip() or file_path.name
            report["imported"].append({"filename": file_path.name, "name": name, "provider": provider})
        except Exception as e:
            report["failed"].append({"filename": file_path.name, "error": str(e)})

    return report
