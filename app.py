import streamlit as st

# CSS - Enlever le vert horrible
st.markdown("""
<style>
    .stButton > button { background-color: #0284c7 !important; color: white !important; }
    .stMetricValue { color: #0284c7 !important; }
</style>
""", unsafe_allow_html=True)
import os, re, json, subprocess, shutil, time, base64
import urllib.request
import pandas as pd
import urllib.parse
import unicodedata
from datetime import datetime
from pptx import Presentation

# ---------------------------------------------------------------------------
# 1. CONFIGURATION ET CHEMINS
# ---------------------------------------------------------------------------
NO_DATA = "---"

# Detection cloud vs local
try:
    IS_CLOUD = "gcp_service_account" in st.secrets
except:
    IS_CLOUD = False

# Chemins
PATH_LOCAL = "/Users/juliensac/Library/CloudStorage/GoogleDrive-julien@miint.pro/Drive partagés/Noovee - CV"

if IS_CLOUD:
    os.makedirs("/tmp/noovee/PDF", exist_ok=True)
    PATH_DOSSIER = "/tmp/noovee"
    PATH_PDF     = "/tmp/noovee/PDF"
    PATH_DB      = "/tmp/noovee/contacts_db.json"
else:
    PATH_DOSSIER = PATH_LOCAL
    PATH_PDF     = os.path.join(PATH_LOCAL, "PDF")
    PATH_DB      = os.path.join(PATH_LOCAL, "contacts_db.json")

ANTHROPIC_URL = "https://api.anthropic.com/v1/messages"
GROQ_URL      = "https://api.groq.com/openai/v1/chat/completions"
GEMINI_URL    = "https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash:generateContent"

MISTRAL_URL = "https://api.mistral.ai/v1/chat/completions"
OPENAI_URL  = "https://api.openai.com/v1/chat/completions"

def charger_toutes_les_cles():
    cles = {}
    env_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), ".env")
    vars_ok = ["ANTHROPIC_API_KEY","MISTRAL_API_KEY","OPENAI_API_KEY","GROQ_API_KEY","GEMINI_API_KEY"]
    for var in vars_ok:
        val = os.environ.get(var, "")
        if val: cles[var] = val
    if os.path.exists(env_path):
        with open(env_path) as f:
            for line in f:
                line = line.strip()
                for var in vars_ok:
                    if line.startswith(var + "=") and var not in cles:
                        cles[var] = line.split("=",1)[1].strip().strip('"').strip("'")
    return cles

ALL_KEYS = charger_toutes_les_cles()

def get_provider_extraction():
    for v in ["MISTRAL_API_KEY","ANTHROPIC_API_KEY","OPENAI_API_KEY","GEMINI_API_KEY"]:
        if v in ALL_KEYS: return v.replace("_API_KEY",""), ALL_KEYS[v]
    return None, ""

def get_provider_scoring():
    for v in ["OPENAI_API_KEY","ANTHROPIC_API_KEY","MISTRAL_API_KEY","GEMINI_API_KEY"]:
        if v in ALL_KEYS: return v.replace("_API_KEY",""), ALL_KEYS[v]
    return None, ""

def get_provider_email():
    for v in ["ANTHROPIC_API_KEY","OPENAI_API_KEY","MISTRAL_API_KEY"]:
        if v in ALL_KEYS: return v.replace("_API_KEY",""), ALL_KEYS[v]
    return None, ""

def appeler_ia(prompt, provider, key, max_tokens=400, model_override=None):
    if not provider or not key: raise Exception("Pas de cle IA")
    if provider == "ANTHROPIC":
        model = model_override or "claude-haiku-4-5-20251001"
        payload = json.dumps({"model":model,"max_tokens":max_tokens,
            "messages":[{"role":"user","content":prompt}]}).encode()
        req = urllib.request.Request(ANTHROPIC_URL, data=payload,
            headers={"Content-Type":"application/json","x-api-key":key,
                     "anthropic-version":"2023-06-01"}, method="POST")
        with urllib.request.urlopen(req, timeout=25) as r:
            return json.loads(r.read())["content"][0]["text"].strip()
    elif provider in ("MISTRAL","OPENAI","GROQ"):
        url   = {"MISTRAL":MISTRAL_URL,"OPENAI":OPENAI_URL,"GROQ":GROQ_URL}[provider]
        model = model_override or {"MISTRAL":"mistral-small-latest","OPENAI":"gpt-4o-mini","GROQ":"llama3-8b-8192"}[provider]
        payload = json.dumps({"model":model,"max_tokens":max_tokens,
            "messages":[{"role":"user","content":prompt}],"temperature":0}).encode()
        req = urllib.request.Request(url, data=payload,
            headers={"Content-Type":"application/json","Authorization":"Bearer "+key}, method="POST")
        with urllib.request.urlopen(req, timeout=25) as r:
            return json.loads(r.read())["choices"][0]["message"]["content"].strip()
    elif provider == "GEMINI":
        payload = json.dumps({"contents":[{"parts":[{"text":prompt}]}]}).encode()
        req = urllib.request.Request(GEMINI_URL+"?key="+key, data=payload,
            headers={"Content-Type":"application/json"}, method="POST")
        with urllib.request.urlopen(req, timeout=25) as r:
            return json.loads(r.read())["candidates"][0]["content"]["parts"][0]["text"].strip()
    raise Exception("Provider inconnu: "+str(provider))

# Compat ancien code
_first = list(ALL_KEYS.items())[0] if ALL_KEYS else (None, "")
AI_PROVIDER = _first[0]
AI_KEY      = _first[1]

def charger_api_key():
    return AI_PROVIDER, AI_KEY


# ---------------------------------------------------------------------------
# 2. NORMALISATION (definie en premier car utilisee partout)
# ---------------------------------------------------------------------------
def normaliser(texte):
    if not texte:
        return ""
    t = unicodedata.normalize("NFD", texte.lower())
    t = "".join(c for c in t if unicodedata.category(c) != "Mn")
    for c in ["\u2019", "\u2018", "\u2032", "'", "`"]:
        t = t.replace(c, " ")
    t = t.replace("-", " ").replace("_", " ")
    t = re.sub(r"[^\w\s]", " ", t)
    return re.sub(r"\s+", " ", t).strip()


# ---------------------------------------------------------------------------
# 3. TABLE DE CORRESPONDANCE ENTREPRISES / SECTEURS
# ---------------------------------------------------------------------------
SECTEURS = {
    "Assurance": [
        "allianz", "axa", "generali", "groupama", "maif", "maaf", "mma", "covea",
        "april", "swisslife", "cardif", "predica", "ag2r", "malakoff", "humanis",
        "macif", "matmut", "sogessur", "aig", "zurich", "chubb", "hiscox",
        "lloyd", "scor", "mutuelle", "prevoyance", "assurance", "assureur",
        "souscription", "sinistre", "iard", "vie prevoyance"
    ],
    "Banque": [
        "bnp paribas", "bnp", "societe generale", "credit agricole", "lcl",
        "banque populaire", "caisse epargne", "bpce", "la banque postale",
        "credit mutuel", "cic", "hsbc", "natixis", "ing", "boursorama",
        "cacib", "credit agricole cib", "ca cib", "sgcib", "sg cib",
        "oddo", "rothschild", "lazard", "ubs", "deutsche bank",
        "jp morgan", "goldman sachs", "morgan stanley", "barclays", "citi",
        "bred", "credit du nord", "banque de france", "banque centrale",
        "banque", "bancaire", "financement", "tresorerie", "back office bancaire",
        "front office", "salle des marches", "marches financiers", "trading"
    ],
    "Conseil": [
        "accenture", "capgemini", "mckinsey", "bcg", "bain", "deloitte", "pwc",
        "kpmg", "ey", "ernst young", "oliver wyman", "roland berger", "sia partners",
        "wavestone", "sopra steria", "cgi", "atos", "ibm consulting", "boston consulting",
        "arthur d little", "altran", "alten", "aubay", "onepoint", "eleven strategy",
        "consulting", "conseil", "management consulting"
    ],
    "Luxe / Cosmetique": [
        "loreal", "l oreal", "lvmh", "chanel", "hermes", "kering", "richemont",
        "dior", "givenchy", "gucci", "saint laurent", "bottega veneta",
        "cartier", "van cleef", "bulgari", "tiffany", "sephora", "lancome",
        "yves saint laurent", "ysl", "clarins", "shiseido", "estee lauder",
        "sisley", "coty", "interparfums", "beauty", "cosmetique", "parfum",
        "haute couture", "maroquinerie", "joaillerie"
    ],
    "Retail / Distribution": [
        "carrefour", "leclerc", "auchan", "casino", "lidl", "aldi", "intermarche",
        "systeme u", "monoprix", "fnac", "darty", "decathlon", "leroy merlin",
        "ikea", "zara", "h&m", "uniqlo", "primark", "kiabi", "la redoute",
        "amazon", "cdiscount", "boulanger", "cultura", "maisons du monde",
        "grande distribution", "retail", "e-commerce", "marketplace"
    ],
    "Telecom / Tech": [
        "orange", "sfr", "bouygues telecom", "free", "numericable", "iliad",
        "microsoft", "google", "apple", "meta", "salesforce", "sap", "oracle",
        "adobe", "servicenow", "workday", "veeva", "dassault systemes",
        "murex", "temenos", "finastra", "digital", "informatique", "logiciel",
        "software", "cloud", "saas", "cybersecurite", "intelligence artificielle"
    ],
    "Energie": [
        "edf", "engie", "totalenergies", "total", "bp", "shell", "equinor",
        "rte", "enedis", "grdf", "grtgaz", "direct energie", "vattenfall",
        "schneider electric", "siemens energy", "vestas", "orsted",
        "energie", "electricite", "gaz naturel", "renouvelable", "solaire", "eolien"
    ],
    "Industrie / Automobile": [
        "renault", "stellantis", "peugeot", "citroen", "ds", "michelin",
        "valeo", "faurecia", "plastic omnium", "saint gobain", "lafarge",
        "holcim", "airbus", "safran", "thales", "dassault aviation",
        "naval group", "nexter", "bouygues", "vinci", "eiffage", "colas",
        "automobile", "aeronautique", "defense", "construction", "btp"
    ],
    "Sante / Pharma": [
        "sanofi", "pierre fabre", "ipsen", "servier", "biomerieux", "essilor",
        "roche", "novartis", "pfizer", "astrazeneca", "johnson", "medtronic",
        "stryker", "abbott", "becton dickinson", "aphp", "ap hp",
        "clinique", "korian", "orpea", "ramsay", "elsan",
        "sante", "pharma", "medicament", "hopital", "clinique", "medical"
    ],
    "Transport / Logistique": [
        "air france", "sncf", "ratp", "transdev", "keolis", "msc", "cma cgm",
        "bollore", "geodis", "dhl", "fedex", "ups", "chronopost", "la poste",
        "uber", "blablacar", "getlink", "eurostar", "thalys",
        "transport", "logistique", "supply chain", "fret", "livraison"
    ],
    "Immobilier": [
        "nexity", "bouygues immobilier", "vinci immobilier", "cogedim",
        "kaufman broad", "foncia", "orpi", "century 21", "laforet",
        "unibail", "klepierre", "icade", "gecina", "covivio", "altarea",
        "emeria", "emeria technologies", "firstime", "immogroup",
        "immobilier", "foncier", "promotion immobiliere", "property", "real estate",
        "gestion locative", "syndic", "copropriete", "bail", "loyer"
    ],
    "Medias / Communication": [
        "tf1", "m6", "france television", "france tv", "canal", "bein sports",
        "lagardere", "hachette", "figaro", "le monde", "liberation", "les echos",
        "publicis", "havas", "ogilvy", "mccann", "bbdo", "dentsu",
        "communication", "marketing", "publicite", "media", "presse", "audiovisuel"
    ],
    "Ressources Humaines": [
        "adecco", "manpower", "randstad", "michael page", "hays", "robert half",
        "apec", "france travail", "pole emploi", "linkedin", "indeed",
        "welcome to the jungle", "edenred", "sodexo",
        "ressources humaines", "rh", "recrutement", "talent", "paie", "sirh"
    ],
    "Agroalimentaire": [
        "danone", "lactalis", "bel", "bongrain", "savencia", "andros",
        "bonduelle", "marie", "fleury michon", "herta", "william saurin",
        "nestle", "unilever", "mars", "ferrero", "lindt", "pernod ricard",
        "moet hennessy", "bacardi", "sodebo", "agro", "agroalimentaire", "alimentaire"
    ],
}

# Charger les secteurs custom s'ils existent (modifies par l'utilisateur)
PATH_SECTEURS_CUSTOM = os.path.join(os.path.dirname(os.path.abspath(__file__)), "secteurs_custom.json")

def charger_secteurs_actifs():
    """Charge secteurs_custom.json si existe, sinon SECTEURS par defaut."""
    if os.path.exists(PATH_SECTEURS_CUSTOM):
        try:
            with open(PATH_SECTEURS_CUSTOM, "r", encoding="utf-8") as f:
                custom = json.load(f)
            # Fusionner : custom en priorite, completer avec SECTEURS pour les manquants
            fusionnes = dict(SECTEURS)
            fusionnes.update(custom)
            return fusionnes
        except Exception:
            pass
    return dict(SECTEURS)

SECTEURS_ACTIFS = charger_secteurs_actifs()

# ---------------------------------------------------------------------------
# DOMAINES FONCTIONNELS
# ---------------------------------------------------------------------------
DOMAINES_FONCTIONNELS = [
    "Conformité / RGPD",
    "Finance / Contrôle de gestion",
    "Achats / Procurement",
    "Ressources Humaines",
    "Marketing / Communication",
    "Commercial / Business Development",
    "IT / Digital / Tech",
    "Data / BI / Analytics",
    "Juridique / Legal",
    "Risk Management / Audit",
    "Supply Chain / Logistique",
    "Transformation / Change Management",
    "Project Management / PMO",
    "Stratégie / Conseil",
    "Comptabilité / Trésorerie",
    "Cybersécurité",
    "Immobilier / Facilities",
    "Communication / Relations Presse",
    "Développement Durable / RSE",
]

# Mots-cles par domaine pour detection automatique
DOMAINES_MOTS_CLES = {
    "Conformité / RGPD":              ["rgpd","conformite","dpo","cnil","pia","aipd","compliance","reglementaire"],
    "Finance / Contrôle de gestion":  ["finance","controle de gestion","tresorerie","budget","comptabilite","ifrs","consolidation"],
    "Achats / Procurement":           ["achats","procurement","sourcing","fournisseur","appel offres","negociation achat"],
    "Ressources Humaines":            ["rh","ressources humaines","recrutement","paie","sirh","gestion talents","formation"],
    "Marketing / Communication":      ["marketing","communication","branding","campagne","digital marketing","content"],
    "Commercial / Business Development": ["commercial","business development","vente","grands comptes","crm","chiffre affaires"],
    "IT / Digital / Tech":            ["informatique","it","digital","developpement","cloud","infrastructure","devops","software"],
    "Data / BI / Analytics":          ["data","bi","business intelligence","analytics","power bi","tableau","sql","python"],
    "Juridique / Legal":              ["juridique","legal","droit","contrats","avocat","contentieux","propriete intellectuelle"],
    "Risk Management / Audit":        ["risk","risque","audit","controle interne","sox","cartographie risques"],
    "Supply Chain / Logistique":      ["supply chain","logistique","entrepot","transport","stock","approvisionnement"],
    "Transformation / Change Management": ["transformation","change management","conduite changement","organisation","restructuration"],
    "Project Management / PMO":       ["chef de projet","pmo","gestion projet","agile","scrum","planification","deliverable"],
    "Stratégie / Conseil":            ["strategie","conseil","consulting","business strategy","plan strategique"],
    "Comptabilité / Trésorerie":      ["comptabilite","comptable","tresorerie","cash flow","bilan","fiscal"],
    "Cybersécurité":                  ["cybersecurite","securite information","ssi","iso 27001","pentest","soc"],
    "Immobilier / Facilities":        ["immobilier","facilities","baux","property management","foncier"],
    "Communication / Relations Presse": ["relations presse","pr","medias","porte parole","attachee presse"],
    "Développement Durable / RSE":    ["rse","developpement durable","esg","environnement","bilan carbone"],
}

def detecter_domaines(texte, competences=None):
    """Detecte les domaines fonctionnels depuis le texte CV et les competences."""
    if not texte and not competences:
        return []
    hay = normaliser((texte or "") + " " + " ".join(competences or []))
    scores_dom = {}
    for domaine, mots in DOMAINES_MOTS_CLES.items():
        score = sum(1 for m in mots if normaliser(m) in hay)
        if score > 0:
            scores_dom[domaine] = score
    seuil = max(scores_dom.values()) * 0.3 if scores_dom else 0
    return [d for d, sc in sorted(scores_dom.items(), key=lambda x: -x[1]) if sc >= seuil][:3]


# Table plate normalisee pour recherche rapide : {mot_cle_normalise: secteur}
TABLE_RAPIDE = {}
for secteur, mots in SECTEURS_ACTIFS.items():
    for mot in mots:
        TABLE_RAPIDE[normaliser(mot)] = secteur


def deduire_secteurs_depuis_table(texte):
    """Detecte tous les secteurs presents dans un CV (pas juste le dominant)."""
    if not texte:
        return []
    t = normaliser(texte)
    scores = {}
    for mot_norme, secteur in TABLE_RAPIDE.items():
        if mot_norme in t:
            poids = 3 if len(mot_norme.split()) >= 2 else 1
            scores[secteur] = scores.get(secteur, 0) + poids
    if not scores:
        return []
    seuil = max(scores.values()) * 0.3
    return [s for s, sc in sorted(scores.items(), key=lambda x: -x[1]) if sc >= seuil]


def detecter_secteur(texte):
    """
    Detecte le secteur depuis le texte du CV.
    Les noms d entreprises (> 4 mots) ont plus de poids que les mots generiques.
    """
    if not texte:
        return NO_DATA
    t = normaliser(texte)
    scores = {}
    for mot_norme, secteur in TABLE_RAPIDE.items():
        if mot_norme in t:
            # Poids : nom compose (ex "credit agricole cib") = 3, mot simple = 1
            poids = 3 if len(mot_norme.split()) >= 2 else 1
            scores[secteur] = scores.get(secteur, 0) + poids
    if scores:
        return max(scores, key=scores.get)
    return NO_DATA


# ---------------------------------------------------------------------------
# 4. EXTRACTION DE TEXTE (PPTX ET PDF)
# ---------------------------------------------------------------------------
def extraire_texte_pptx(path):
    try:
        prs = Presentation(path)
        texte = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texte += " " + shape.text
        return texte.strip() or None
    except Exception:
        return None

def extraire_texte_pdf(path):
    try:
        import pypdf
        texte = ""
        with open(path, "rb") as f:
            reader = pypdf.PdfReader(f)
            for page in reader.pages:
                texte += " " + (page.extract_text() or "")
        return texte.strip() or None
    except Exception:
        return None


# ---------------------------------------------------------------------------
# 5. EXTRACTION IA (CLAUDE HAIKU EN PRIORITE)
# ---------------------------------------------------------------------------
def analyser_cv_avec_ia(texte):
    """
    Extraction CV enrichie v2 :
    - nom, prenom, email, tel, poste, competences
    - entreprises avec secteur
    - domaines fonctionnels
    - 2-3 experiences cles
    Provider : Mistral > Claude > OpenAI > Gemini
    """
    if not texte: return None
    provider, key = get_provider_extraction()
    if not key: return None

    domaines_liste = ", ".join(DOMAINES_FONCTIONNELS[:10])

    prompt = "\n".join([
        "Analyse ce CV et extrais les informations en JSON uniquement, sans markdown.",
        "Format exact :",
        (
            '{"nom":"NOM","prenom":"Prenom","email":"email@ex.com","telephone":"0600000000",'
            '"poste":"Poste principal","annees_experience":5,'
            '"competences":["comp1","comp2","comp3","comp4","comp5"],'
            '"domaines_fonctionnels":["domaine1","domaine2"],'
            '"entreprises":[{"nom":"Entreprise","secteur":"Secteur","annees":2}],'
            '"experiences_cles":['
            '{"titre":"Poste - Entreprise (2022-2024)",'
            '"description":"Resume en 1 phrase des missions et realisations cles"}'
            ']}'
        ),
        "Domaines fonctionnels possibles : " + domaines_liste,
        "Extrais 2 ou 3 experiences les plus significatives (pas toutes).",
        "Telephone commence par 06 ou 07. Si info absente mets null.",
        "TEXTE COMPLET DU CV :", texte[:6000]
    ])
    try:
        raw = appeler_ia(prompt, provider, key, max_tokens=700)
        raw = raw.replace("```json","").replace("```","").strip()
        return json.loads(raw)
    except Exception:
        return None


# ---------------------------------------------------------------------------
# 6. GESTION DE LA BASE DE DONNEES
# ---------------------------------------------------------------------------
def charger_db():
    return charger_db_cloud()

def sauvegarder_db(db):
    sauvegarder_db_cloud(db)

def tel_valide(tel):
    return bool(tel) and str(tel) not in (NO_DATA, "", "None", "null") and len(str(tel)) >= 10

def lien_whatsapp(telephone, message):
    return "https://web.whatsapp.com/send?phone=33" + telephone[1:] + "&text=" + urllib.parse.quote(message)


# ---------------------------------------------------------------------------
# 7. CHARGEMENT DES CVS (AVEC CACHE)
# ---------------------------------------------------------------------------
@st.cache_data(show_spinner="Chargement des CVs...")
def charger_tous_les_cvs(path_dossier, path_pdf):
    entrees = {}

    # Si cloud : telecharger les fichiers depuis Drive
    if IS_CLOUD:
        folder_id, pdf_folder_id = drive_get_folder_ids()
        if folder_id:
            # Telecharger PPTX
            for f in drive_list_files(folder_id, (".pptx",)):
                dest = os.path.join(path_dossier, f["name"])
                if not os.path.exists(dest):
                    drive_download_file(f["id"], dest)
            # Telecharger PDF
            if pdf_folder_id:
                for f in drive_list_files(pdf_folder_id, (".pdf",)):
                    dest = os.path.join(path_pdf, f["name"])
                    if not os.path.exists(dest):
                        drive_download_file(f["id"], dest)

    def ajouter(base, texte, source, pdf_path):
        if base not in entrees:
            entrees[base] = {"texte": texte or "", "source": source, "pdf_path": pdf_path}
        else:
            if not entrees[base]["texte"] and texte:
                entrees[base]["texte"] = texte
            if not entrees[base]["pdf_path"] and pdf_path:
                entrees[base]["pdf_path"] = pdf_path

    for f in sorted(os.listdir(path_dossier)):
        if not f.endswith(".pptx"):
            continue
        base = f.replace(".pptx", "")

        # Cherche le PDF en priorite (meilleure qualite d extraction)
        pdf_path = None
        for c in [os.path.join(path_pdf, base + ".pdf"), os.path.join(path_dossier, base + ".pdf")]:
            if os.path.exists(c):
                pdf_path = c
                break

        # PDF en priorite, PPTX en fallback
        if pdf_path:
            texte = extraire_texte_pdf(pdf_path)
            source = "pdf"
            if not texte:
                texte = extraire_texte_pptx(os.path.join(path_dossier, f))
                source = "pptx"
        else:
            texte = extraire_texte_pptx(os.path.join(path_dossier, f))
            source = "pptx" if texte else "illisible"

        ajouter(base, texte, source, pdf_path)

    if os.path.exists(path_pdf):
        for f in sorted(os.listdir(path_pdf)):
            if not f.endswith(".pdf"):
                continue
            base = f.replace(".pdf", "")
            pdf_path = os.path.join(path_pdf, f)
            texte = None if (base in entrees and entrees[base]["texte"]) else extraire_texte_pdf(pdf_path)
            ajouter(base, texte, "pdf", pdf_path)

    for f in sorted(os.listdir(path_dossier)):
        if not f.endswith(".pdf"):
            continue
        base = f.replace(".pdf", "")
        pdf_path = os.path.join(path_dossier, f)
        texte = None if (base in entrees and entrees[base]["texte"]) else extraire_texte_pdf(pdf_path)
        ajouter(base, texte, "pdf", pdf_path)

    return entrees


# ---------------------------------------------------------------------------
# 8. RECHERCHE
# ---------------------------------------------------------------------------
def filtrer_cvs(entrees, query, db):
    termes = normaliser(query).split()
    if not termes:
        return []
    resultats = []
    for base, data in entrees.items():
        info = db.get(base, {})
        secteur = info.get("secteur", NO_DATA)

        # Enrichir le haystack avec les synonymes des secteurs mentionnes
        # Ex: recherche "immobilier" → ajoute tous les mots-cles du secteur Immobilier
        extra_secteurs = ""
        for terme in termes:
            for sect_nom, sect_mots in SECTEURS_ACTIFS.items():
                if terme in normaliser(sect_nom):
                    extra_secteurs += " " + " ".join(sect_mots)
                    # Aussi si le terme est un mot-cle du secteur → ajoute le nom du secteur
                elif any(terme in normaliser(m) for m in sect_mots):
                    extra_secteurs += " " + sect_nom

        # Multi-secteurs + entreprises dans le haystack
        secteurs_list    = info.get("secteurs", [secteur])
        domaines_list    = info.get("domaines", [])
        entreprises_list = " ".join(
            e.get("nom","") + " " + e.get("secteur","")
            for e in info.get("entreprises", [])
        )
        exps_list = " ".join(
            e.get("titre","") + " " + e.get("description","")
            for e in info.get("experiences_cles", [])
        )
        haystack = normaliser(
            base + " " +
            data.get("texte", "") + " " +
            info.get("nom", "") + " " +
            info.get("prenom", "") + " " +
            info.get("email", "") + " " +
            info.get("poste", "") + " " +
            " ".join(secteurs_list) + " " +
            " ".join(domaines_list) + " " +
            " ".join(info.get("competences", [])) + " " +
            entreprises_list + " " +
            exps_list + " " +
            extra_secteurs
        )
        if not all(t in haystack for t in termes):
            continue
        tel = info.get("telephone", NO_DATA)
        resultats.append({
            "Selec.":      False,
            "Nom":         info.get("nom", NO_DATA),
            "Prenom":      info.get("prenom", NO_DATA),
            "Telephone":   tel,
            "Secteur":     secteur,
            "Poste":       info.get("poste", NO_DATA),
            "Competences": ", ".join(info.get("competences", [])),
            "Fichier":     base,
            "_pdf_path":   data["pdf_path"],
        })
    return resultats


# ---------------------------------------------------------------------------
# 9. INTERFACE STREAMLIT
# ---------------------------------------------------------------------------
st.set_page_config(page_title="Noovee - Contacts IA", layout="wide", page_icon="🌿")

# CSS Noovee — design system complet
st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Syne:wght@600;700;800&family=DM+Sans:ital,wght@0,300;0,400;0,500;0,600;1,400&display=swap');

/* ─── Variables ────────────────────────────── */
:root {
    --g-dark:   #093F28;
    --g-mid:    #1A6B45;
    --g-light:  #E8F4EE;
    --coral:    #E07878;
    --coral-lt: #FDF1F1;
    --navy:     #1E2A3A;
    --slate:    #486074;
    --muted:    #8898AA;
    --bg:       #F5F7FA;
    --white:    #FFFFFF;
    --border:   #E2E8F0;
    --shadow:   0 2px 12px rgba(9,63,40,.10);
    --radius:   12px;
    --radius-sm:8px;
}

/* ─── Force theme clair ────────────────────── */
html, body,
[data-testid="stAppViewContainer"],
[data-testid="stMain"],
section[data-testid="stSidebar"],
.stApp { 
    background: var(--bg) !important;
    background-color: var(--bg) !important;
    color: var(--navy) !important;
    font-family: 'DM Sans', sans-serif !important;
    color-scheme: light !important;
}
[data-testid="stSidebar"] { display: none !important; }
[data-testid="stToolbar"]  { display: none !important; }

/* Forcer tous les textes en sombre */
p, span, div, label, h1, h2, h3, h4 {
    color: var(--navy) !important;
}

/* Corriger boutons noirs (mode sombre system) */
div.stButton > button {
    background-color: var(--white) !important;
    color: var(--navy) !important;
}
div.stButton > button[kind="primary"] {
    background-color: var(--g-dark) !important;
    color: white !important;
}

/* Corriger fond des containers */
[data-testid="stVerticalBlock"],
[data-testid="stHorizontalBlock"] {
    background: transparent !important;
}

/* Retirer padding excessif */
.main .block-container {
    padding-top: 1.5rem !important;
    padding-bottom: 3rem !important;
    max-width: 1280px !important;
}

/* ─── Header ───────────────────────────────── */
.noovee-header {
    background: linear-gradient(135deg, var(--g-dark) 0%, var(--g-mid) 100%);
    padding: 1.4rem 2rem;
    border-radius: var(--radius);
    margin-bottom: 1.5rem;
    display: flex;
    align-items: center;
    gap: 1.2rem;
    box-shadow: var(--shadow);
    position: relative;
    overflow: hidden;
}
.noovee-header::after {
    content: '';
    position: absolute;
    right: -40px; top: -40px;
    width: 200px; height: 200px;
    border-radius: 50%;
    background: rgba(255,255,255,.05);
}
.noovee-logo {
    width: 48px; height: 48px;
    background: white;
    border-radius: 50%;
    display: flex; align-items: center; justify-content: center;
    font-family: 'Syne', sans-serif;
    font-size: 1.5rem; font-weight: 800;
    color: var(--g-dark);
    flex-shrink: 0;
    box-shadow: 0 2px 8px rgba(0,0,0,.15);
}
.noovee-header h1 {
    font-family: 'Syne', sans-serif !important;
    color: white !important;
    font-size: 1.5rem !important;
    font-weight: 700 !important;
    margin: 0 0 2px 0 !important;
    letter-spacing: 0.3px;
}
.noovee-header .subtitle {
    color: rgba(255,255,255,.70);
    font-size: 0.82rem;
    margin: 0;
    font-weight: 300;
}

/* ─── Tabs ─────────────────────────────────── */
[data-testid="stTabs"] [data-baseweb="tab-list"] {
    background: var(--white) !important;
    border-radius: var(--radius) var(--radius) 0 0 !important;
    padding: 0 1rem !important;
    border-bottom: 2px solid var(--border) !important;
    gap: 0 !important;
}
[data-testid="stTabs"] button[data-baseweb="tab"] {
    font-family: 'DM Sans', sans-serif !important;
    font-weight: 500 !important;
    font-size: 0.88rem !important;
    color: var(--slate) !important;
    padding: 0.9rem 1.2rem !important;
    border-radius: 0 !important;
    border-bottom: 2px solid transparent !important;
    margin-bottom: -2px !important;
    transition: all .2s !important;
}
[data-testid="stTabs"] button[data-baseweb="tab"]:hover {
    color: var(--g-mid) !important;
    background: var(--g-light) !important;
}
[data-testid="stTabs"] button[aria-selected="true"] {
    color: var(--g-dark) !important;
    border-bottom-color: var(--g-dark) !important;
    font-weight: 600 !important;
}

/* ─── Boutons ──────────────────────────────── */
div.stButton > button {
    font-family: 'DM Sans', sans-serif !important;
    font-weight: 500 !important;
    border-radius: var(--radius-sm) !important;
    transition: all .2s !important;
    border: 1.5px solid var(--border) !important;
    font-size: 0.875rem !important;
}
div.stButton > button:hover {
    transform: translateY(-1px) !important;
    box-shadow: 0 4px 12px rgba(0,0,0,.1) !important;
}
div.stButton > button[kind="primary"] {
    background: var(--g-dark) !important;
    border-color: var(--g-dark) !important;
    color: white !important;
    font-weight: 600 !important;
}
div.stButton > button[kind="primary"]:hover {
    background: var(--g-mid) !important;
    border-color: var(--g-mid) !important;
}

/* ─── Inputs ───────────────────────────────── */
div[data-baseweb="input"] input,
div[data-baseweb="textarea"] textarea {
    font-family: 'DM Sans', sans-serif !important;
    border-radius: var(--radius-sm) !important;
    border-color: var(--border) !important;
    background: var(--white) !important;
    font-size: 0.9rem !important;
    transition: border-color .2s !important;
}
div[data-baseweb="input"] input:focus,
div[data-baseweb="textarea"] textarea:focus {
    border-color: var(--g-mid) !important;
    box-shadow: 0 0 0 3px rgba(26,107,69,.12) !important;
}

/* ─── Métriques ────────────────────────────── */
div[data-testid="metric-container"] {
    background: var(--white);
    border-radius: var(--radius);
    padding: 1rem 1.2rem !important;
    border: 1px solid var(--border);
    border-left: 4px solid var(--g-dark) !important;
    box-shadow: var(--shadow);
}
div[data-testid="metric-container"] label {
    font-family: 'DM Sans', sans-serif !important;
    font-size: 0.78rem !important;
    font-weight: 500 !important;
    color: var(--muted) !important;
    text-transform: uppercase;
    letter-spacing: 0.5px;
}
div[data-testid="metric-container"] [data-testid="stMetricValue"] {
    font-family: 'Syne', sans-serif !important;
    font-size: 2rem !important;
    font-weight: 700 !important;
    color: var(--g-dark) !important;
}

/* ─── Cards sections ───────────────────────── */
.card {
    background: var(--white);
    border-radius: var(--radius);
    border: 1px solid var(--border);
    padding: 1.2rem 1.4rem;
    box-shadow: var(--shadow);
    margin-bottom: 1rem;
}
.section-label {
    display: flex;
    align-items: center;
    gap: 0.5rem;
    padding: 0.7rem 1rem;
    border-radius: var(--radius-sm);
    margin-bottom: 0.8rem;
    font-weight: 600;
    font-size: 0.9rem;
}
.section-green {
    background: var(--g-light);
    border-left: 4px solid var(--g-dark);
    color: var(--g-dark);
}
.section-ao {
    background: #EEF2FF;
    border-left: 4px solid #4F46E5;
    color: #4F46E5;
}

/* ─── Dataframe ────────────────────────────── */
[data-testid="stDataFrame"] {
    border-radius: var(--radius) !important;
    border: 1px solid var(--border) !important;
    overflow: hidden !important;
}

/* ─── Progress bar ─────────────────────────── */
[data-testid="stProgress"] > div > div {
    background: var(--g-mid) !important;
    border-radius: 4px !important;
}

/* ─── Alerts ───────────────────────────────── */
[data-testid="stAlert"] {
    border-radius: var(--radius-sm) !important;
    font-family: 'DM Sans', sans-serif !important;
}

/* ─── Expander ─────────────────────────────── */
[data-testid="stExpander"] {
    border: 1px solid var(--border) !important;
    border-radius: var(--radius-sm) !important;
    background: var(--white) !important;
}

/* ─── Divider ──────────────────────────────── */
hr {
    border-color: var(--border) !important;
    margin: 1.5rem 0 !important;
}

/* ─── Caption / small text ─────────────────── */
small, [data-testid="stCaptionContainer"] {
    color: var(--muted) !important;
    font-size: 0.8rem !important;
}

/* ─── Selectbox ────────────────────────────── */
[data-baseweb="select"] {
    border-radius: var(--radius-sm) !important;
}

/* ─── File uploader ────────────────────────── */
[data-testid="stFileUploader"] {
    border-radius: var(--radius-sm) !important;
}
[data-testid="stFileUploadDropzone"] {
    border-radius: var(--radius-sm) !important;
    border-color: var(--border) !important;
    background: var(--white) !important;
}

/* ─── Success / warning / error ────────────── */
[data-testid="stAlert"][kind="success"] {
    background: var(--g-light) !important;
    border-color: var(--g-mid) !important;
    color: var(--g-dark) !important;
}

/* ─── Subheader ────────────────────────────── */
h2, h3 {
    font-family: 'Syne', sans-serif !important;
    color: var(--navy) !important;
    font-weight: 700 !important;
}

/* ─── Link buttons ─────────────────────────── */
a[data-testid="stLinkButton"] > button {
    border-radius: var(--radius-sm) !important;
    font-family: 'DM Sans', sans-serif !important;
}
</style>
""", unsafe_allow_html=True)

# Header Noovee
nb_cvs_total = len(entrees_cvs) if 'entrees_cvs' in dir() else 0
st.markdown("""
<div class="noovee-header">
    <div class="noovee-logo">N</div>
    <div style="flex:1">
        <h1>Noovee — Base de Contacts IA</h1>
        <p class="subtitle">Sourcing de profils sur mesure &nbsp;·&nbsp; Votre sélection sous 48h</p>
    </div>
    <div style="text-align:right;color:rgba(255,255,255,.6);font-size:0.78rem;font-family:'DM Sans',sans-serif;">
        Propulsé par<br>
        <span style="color:white;font-weight:600;">Claude Haiku</span>
    </div>
</div>
""", unsafe_allow_html=True)

if not os.path.exists(PATH_DOSSIER):
    st.error("Dossier introuvable : " + PATH_DOSSIER)
    st.stop()

if "ao_cv_ouvert" not in st.session_state:
    st.session_state["ao_cv_ouvert"] = None
if "ao_criteres" not in st.session_state:
    st.session_state["ao_criteres"] = None
if "ao_cdc" not in st.session_state:
    st.session_state["ao_cdc"] = ""
if "ao_client" not in st.session_state:
    st.session_state["ao_client"] = ""

os.makedirs(PATH_PDF, exist_ok=True)
entrees_cvs = charger_tous_les_cvs(PATH_DOSSIER, PATH_PDF)
db = charger_db()

# Enregistrement automatique des nouveaux CVs + detection secteur
modif = False
for base, data in entrees_cvs.items():
    texte = data.get("texte", "")
    if base not in db:
        match_tel = re.search(r"(0[67](?:[\s.\-]*\d{2}){4})", texte or "")
        tel = "".join(re.findall(r"\d", match_tel.group(1))) if match_tel else NO_DATA
        match_email = re.search(r"[a-zA-Z0-9._%+\-]+@[a-zA-Z0-9.\-]+\.[a-zA-Z]{2,}", texte or "")
        email = match_email.group(0).lower() if match_email else NO_DATA
        db[base] = {
            "nom": NO_DATA, "prenom": NO_DATA,
            "email": email, "telephone": tel,
            "poste": NO_DATA, "competences": [],
            "secteur": detecter_secteur(texte),
            "texte_brut": texte[:8000] if texte else "",  # texte complet stocké
            "ia_enrichi": False,
            "date_ajout": datetime.now().strftime("%Y-%m-%d %H:%M"),
            "whatsapp_envoye": False,
        }
        modif = True
    elif not db[base].get("secteur") or db[base].get("secteur") == NO_DATA:
        # Met a jour le secteur pour les contacts existants sans secteur
        db[base]["secteur"] = detecter_secteur(texte)
        modif = True

if modif:
    sauvegarder_db(db)

# Recalcul automatique secteurs + domaines + experiences synthétiques au démarrage
recalcules = 0
for base, data in entrees_cvs.items():
    if base not in db:
        continue
    texte = data.get("texte", "") or ""
    comps = db[base].get("competences", [])
    t = normaliser(texte)

    # Recalcul secteurs
    scores = {}
    for mot_norme, secteur in TABLE_RAPIDE.items():
        if mot_norme in t:
            poids = 3 if len(mot_norme.split()) >= 2 else 1
            scores[secteur] = scores.get(secteur, 0) + poids
    if scores:
        nouveau_secteur = max(scores, key=scores.get)
        if db[base].get("secteur") != nouveau_secteur:
            db[base]["secteur"] = nouveau_secteur
            recalcules += 1

    # Mettre à jour le texte brut si absent ou vide
    if not db[base].get("texte_brut") and texte:
        db[base]["texte_brut"] = texte[:8000]
        recalcules += 1

    # Domaines fonctionnels
    if not db[base].get("domaines"):
        domaines_auto = detecter_domaines(texte, comps)
        if domaines_auto:
            db[base]["domaines"] = domaines_auto
            recalcules += 1

    # Normaliser les competences trop longues (> 4 mots → extraire termes courts)
    comps_actuelles = db[base].get("competences", [])
    comps_normalisees = []
    for c in comps_actuelles:
        mots = c.split()
        if len(mots) <= 3:
            comps_normalisees.append(c)  # deja court
        else:
            # Garder quand meme mais ajouter les mots-cles importants extraits
            comps_normalisees.append(c)
            # Extraire les termes techniques courts contenus dedans
            termes_courts = ["rgpd","dpo","pia","dpia","ropa","cnil","gdpr",
                             "ifrs","sql","python","agile","scrum","power bi",
                             "audit","conformite","risk","sap","erp","crm"]
            for terme in termes_courts:
                if terme in normaliser(c) and terme not in [normaliser(x) for x in comps_normalisees]:
                    comps_normalisees.append(terme.upper())
    if len(comps_normalisees) != len(comps_actuelles):
        db[base]["competences"] = comps_normalisees
        recalcules += 1

    # Générer experiences synthétiques si manquantes mais domaines disponibles
    if not db[base].get("experiences") and db[base].get("domaines"):
        domaines_cv = db[base]["domaines"]
        annees_tot = db[base].get("annees_experience", 0)
        if annees_tot == 0:
            # Estimer depuis les dates dans le texte
            yrs = re.findall(r"20(\d{2})", texte)
            if len(yrs) >= 2:
                annees_tot = max(int(y) for y in yrs) - min(int(y) for y in yrs)
        # Créer une experience synthétique par domaine
        exps_synth = []
        for dom in domaines_cv[:3]:
            mots_dom = DOMAINES_MOTS_CLES.get(dom, [])
            mots_trouves = [m for m in mots_dom if normaliser(m) in t][:5]
            exps_synth.append({
                "poste":       db[base].get("poste", dom),
                "entreprise":  "Parcours professionnel",
                "domaine":     dom,
                "annees":      max(1, annees_tot // max(len(domaines_cv), 1)),
                "mots_cles":   mots_trouves,
                "synthetique": True
            })
        if exps_synth:
            db[base]["experiences"] = exps_synth
            recalcules += 1

if recalcules > 0:
    sauvegarder_db(db)

# ---------------------------------------------------------------------------
# HELPER GLOBAL : afficher un CV PDF inline
# ---------------------------------------------------------------------------
def afficher_cv_inline(pdf_path, nom, cle_unique, session_key=None):
    """Affiche un CV PDF en pleine largeur avec PDF.js. session_key = cle a remettre a None pour fermer."""
    if not pdf_path or not os.path.exists(pdf_path):
        st.warning("PDF introuvable.")
        return
    col_t, col_dl, col_close = st.columns([4, 1, 1])
    with col_t:
        st.markdown(f"**📄 CV — {nom}**")
    with col_dl:
        with open(pdf_path, "rb") as fp:
            st.download_button("⬇️ PDF", data=fp,
                file_name=os.path.basename(pdf_path),
                mime="application/pdf",
                key="dl_inline_" + cle_unique)
    with col_close:
        if st.button("✕ Fermer", key="close_inline_" + cle_unique):
            if session_key:
                st.session_state[session_key] = None
            st.rerun()
    with open(pdf_path, "rb") as fp:
        b64 = base64.b64encode(fp.read()).decode()
    uid = re.sub(r"[^a-zA-Z0-9]", "", cle_unique)[:12]
    html = (
        f"<script src='https://cdnjs.cloudflare.com/ajax/libs/pdf.js/3.11.174/pdf.min.js'></script>"
        f"<div id='cv{uid}' style='background:#f5f5f5;padding:8px;border-radius:10px;"
        f"overflow-y:auto;max-height:820px;border:1px solid #e2e8f0;margin-top:8px'></div>"
        f"<script>"
        f"pdfjsLib.GlobalWorkerOptions.workerSrc='https://cdnjs.cloudflare.com/ajax/libs/pdf.js/3.11.174/pdf.worker.min.js';"
        f"(function(){{"
        f"var b='{b64}';var bn=atob(b);var ar=new Uint8Array(bn.length);"
        f"for(var i=0;i<bn.length;i++) ar[i]=bn.charCodeAt(i);"
        f"pdfjsLib.getDocument({{data:ar}}).promise.then(function(pdf){{"
        f"var c=document.getElementById('cv{uid}');"
        f"for(var p=1;p<=pdf.numPages;p++){{"
        f"(function(pn){{pdf.getPage(pn).then(function(page){{"
        f"var vp=page.getViewport({{scale:1.5}});"
        f"var cv=document.createElement('canvas');"
        f"cv.width=vp.width;cv.height=vp.height;"
        f"cv.style.cssText='width:100%;margin-bottom:6px;background:white;border-radius:4px;display:block;';"
        f"c.appendChild(cv);"
        f"page.render({{canvasContext:cv.getContext('2d'),viewport:vp}});"
        f"}})}})(p);"
        f"}}"
        f"}});"
        f"}})();"
        f"</script>"
    )
    st.components.v1.html(html, height=860, scrolling=True)


# ONGLETS
tab_recherche, tab_contacts, tab_doublons, tab_secteurs, tab_outils = st.tabs([
    "Recherche", "Base de contacts", "Doublons", "Table secteurs", "Statistiques"
])


# ---------------------------------------------------------------------------
# ONGLET 1 — RECHERCHE + MATCHING AO
# ---------------------------------------------------------------------------
with tab_recherche:

    col_gauche, col_droite = st.columns([1, 1], gap="large")

    # ========================
    # COLONNE GAUCHE : RECHERCHE
    # ========================
    with col_gauche:
        st.markdown('<div class="section-label section-green">🔍 &nbsp;Recherche par compétence</div>', unsafe_allow_html=True)

        query = st.text_input("Nom, secteur, compétence, entreprise...", "", key="query_main")
        st.caption(str(len(entrees_cvs)) + " CVs indexés")

    # ========================
    # COLONNE DROITE : MATCHING AO
    # ========================
    with col_droite:
        st.markdown('<div class="section-label section-ao">🎯 Matching Appel d\'Offres</div>', unsafe_allow_html=True)

        nom_client_ao = st.session_state.get("ao_client", "")
        cdc_texte_ao  = st.text_area("Coller le cahier des charges :", height=120,
            placeholder="Collez ici le texte de l appel d offres...", key="cdc_input")

        cdc_file_ao = st.file_uploader("Ou uploader un PDF", type=["pdf"], key="cdc_file_ao")
        if cdc_file_ao:
            try:
                import pypdf, io
                reader = pypdf.PdfReader(io.BytesIO(cdc_file_ao.read()))
                cdc_texte_ao = " ".join(page.extract_text() or "" for page in reader.pages)
                st.success("PDF lu — " + str(len(cdc_texte_ao)) + " caracteres")
            except Exception as e:
                st.error("Erreur PDF : " + str(e))

        if st.button("Analyser et matcher", type="primary",
                     disabled=not cdc_texte_ao, key="btn_ao"):
            if not AI_KEY:
                st.error("Cle IA non configuree.")
            else:
                with st.spinner("Analyse IA du cahier des charges..."):
                    prompt_ao = "\n".join([
                        "Analyse ce cahier des charges et extrais les informations cles en JSON uniquement.",
                        "Format exact sans markdown :",
                        '{"poste":"intitule","competences":["c1","c2","c3","c4","c5"],"secteurs":["s1","s2"],"duree":"6 mois","localisation":"Paris","resume":"2 phrases"}',
                        "Si info absente mets null. CDC :", cdc_texte_ao[:4000]
                    ])
                    try:
                        prov_cdc, key_cdc = get_provider_scoring()
                        raw = appeler_ia(prompt_ao, prov_cdc, key_cdc, max_tokens=400)
                        raw = raw.replace("```json","").replace("```","").strip()
                        criteres_ao = json.loads(raw)
                        st.session_state["ao_criteres"]  = criteres_ao
                        st.session_state["ao_cdc"]       = cdc_texte_ao
                        st.session_state["ao_client"]    = nom_client_ao
                    except Exception as e:
                        st.error("Erreur IA : " + str(e))

        # Affichage resultats matching si dispo
        if st.session_state.get("ao_criteres"):
            crit   = st.session_state["ao_criteres"]
            client = st.session_state.get("ao_client","le client")
            st.success("**" + (crit.get("poste") or "Mission") + "** · " +
                       (crit.get("duree") or "") + " · " + (crit.get("localisation") or ""))
            if crit.get("resume"):
                st.caption(crit["resume"])

            comp_ao    = crit.get("competences") or []
            secteur_ao = crit.get("secteurs") or []
            mots_ao    = [normaliser(c) for c in comp_ao + secteur_ao]

            # Pre-calcul rapide pour tous les profils AVANT d afficher le tableau
            # On applique scorer_criteres a TOUS pour ne pas filtrer Cassandre
            scores_ao = []
            for base, data in entrees_cvs.items():
                info = db.get(base, {})
                nom_s = (info.get("prenom","") + " " + info.get("nom","")).strip()
                if not nom_s or nom_s.strip() in ("---","--- ---"):
                    nom_s = base[:25]
                pdf_s = None
                for cp in [os.path.join(PATH_PDF, base+".pdf"),
                            os.path.join(PATH_DOSSIER, base+".pdf")]:
                    if os.path.exists(cp): pdf_s = cp; break
                scores_ao.append({
                    "base":  base,
                    "nom":   nom_s,
                    "poste": info.get("poste","") or "",
                    "secteur": info.get("secteur","") or "",
                    "tel":   info.get("telephone", NO_DATA),
                    "email": info.get("email", NO_DATA),
                    "pdf":   pdf_s,
                })
            # Appliquer scorer_criteres a tous (pas seulement top5)
            cdc_stock = st.session_state.get("ao_cdc","")

            # ── Helpers scoring ──────────────────────────────
            def nb_etoiles(pct):
                if pct >= 80: return "⭐⭐⭐⭐⭐"
                if pct >= 60: return "⭐⭐⭐⭐"
                if pct >= 40: return "⭐⭐⭐"
                if pct >= 20: return "⭐⭐"
                return "⭐"

            # Dictionnaire d expansion semantique — synonymes metier
            EXPANSION = {
                "rgpd": ["rgpd","gdpr","donnees personnelles","data protection","conformite","dpo",
                         "cnil","privacy","registre traitement","pia","aipd","dpia","accountability",
                         "cookies","ropa","sous traitant","droits personnes"],
                "conformite": ["conformite","compliance","audit","reglementaire","norme","certification",
                               "controle interne","veille reglementaire","gouvernance"],
                "juridique": ["juridique","droit","juriste","legal","avocat","contentieux","contrat"],
                "banque": ["banque","bancaire","credit","finance","assurance","bfin","fintech"],
                "retail": ["retail","grande distribution","commerce","carrefour","leclerc","distribution"],
                "data": ["donnees","data","base de donnees","sql","bi","reporting","power bi","tableau"],
                "projet": ["projet","pmo","pilotage","methodologie","agile","scrum","chef de projet"],
                "dpo": ["dpo","delegue protection","data protection officer","dpd"],
                "pia": ["pia","aipd","dpia","analyse impact","privacy impact"],
            }

            def expanser_termes(mots_ao, cdc_txt):
                """Expanse les termes de l AO avec leurs synonymes metier."""
                termes_exp = set(mots_ao)
                cdc_n = normaliser(cdc_txt[:3000])
                for famille, synonymes in EXPANSION.items():
                    if any(normaliser(s) in cdc_n for s in synonymes[:5]):
                        for s in synonymes:
                            termes_exp.add(normaliser(s))
                return list(termes_exp)

            def scorer_criteres(s_in, cdc_txt, crit_in):
                info        = db.get(s_in["base"], {})
                # Texte complet depuis DB (plus fiable que le cache)
                texte       = (info.get("texte_brut") or
                               entrees_cvs.get(s_in["base"], {}).get("texte", "") or "")
                comps       = info.get("competences", [])
                poste       = s_in.get("poste", "") or ""
                secteurs    = info.get("secteurs", [s_in.get("secteur","")]) or []
                entreprises = info.get("entreprises", [])
                domaines    = info.get("domaines", [])
                exps_cles   = info.get("experiences_cles", [])
                haystack = normaliser(
                    texte + " " + " ".join(comps) + " " + poste + " " +
                    " ".join(secteurs) + " " +
                    " ".join(domaines) + " " +
                    " ".join(e.get("nom","") + " " + e.get("secteur","") for e in entreprises) + " " +
                    " ".join(e.get("titre","") + " " + e.get("description","") for e in exps_cles)
                )

                # ── BONUS DOMAINE FONCTIONNEL (signal fort) ──────────────────
                cdc_n = normaliser(cdc_txt[:3000])
                bonus_domaine = 0

                # Si domaines pas encore calcules → detection auto depuis texte+competences
                domaines_effectifs = domaines if domaines else detecter_domaines(texte, comps)

                for dom in domaines_effectifs:
                    mots_dom = DOMAINES_MOTS_CLES.get(dom, [])
                    if sum(1 for m in mots_dom if normaliser(m) in cdc_n) >= 2:
                        bonus_domaine = 25
                        break
                    elif any(normaliser(m) in cdc_n for m in mots_dom[:3]):
                        bonus_domaine = max(bonus_domaine, 12)

                # Bonus supplementaire si competences directement dans mots-cles CDC
                for comp in comps:
                    comp_n = normaliser(comp)
                    if len(comp_n) > 4 and comp_n in cdc_n:
                        bonus_domaine = min(30, bonus_domaine + 5)

                # ── 1. COMPETENCES FONCTIONNELLES (50%) ──────────────────────
                mots_ao_base = [normaliser(c) for c in (crit_in.get("competences") or [])]
                mots_ao_exp  = expanser_termes(mots_ao_base, cdc_txt)

                # Haystack etendu : texte complet + competences
                comps_n = [normaliser(c) for c in comps]
                poste_n = normaliser(poste)

                # Score competences : dénominateur = uniquement mots_ao_base
                sc_base = 0
                for mc in mots_ao_base:
                    if not mc: continue
                    if any(mc in cn or cn in mc for cn in comps_n if cn):
                        sc_base += 10   # match direct dans compétences DB
                    elif mc in poste_n or any(mc in normaliser(e.get("titre","")) for e in exps_cles):
                        sc_base += 7    # match dans poste ou expériences clés
                    elif mc in haystack:
                        sc_base += 4    # match dans texte brut

                # Score synonymes = BONUS uniquement (ne pénalise pas le dénominateur)
                bonus_exp = 0
                for mc in mots_ao_exp:
                    if not mc or mc in mots_ao_base: continue
                    if any(mc in cn or cn in mc for cn in comps_n if cn): bonus_exp += 3
                    elif mc in haystack: bonus_exp += 1

                # Dénominateur = mots_ao_base × 10 (max possible)
                sc_comp = min(100, int(sc_base / max(len(mots_ao_base) * 10, 1) * 100)
                              + min(25, bonus_exp))

                # ── 2. EXPERIENCE (25%) ──────────────────────────────────────
                annees_cv = 0
                # Detection "X ans d experience" ou "plus de X ans"
                for pat in [r"plus\s+de\s+(\d+)\s*ans?", r"(\d{1,2})\s*ans?\s+d.exp",
                            r"depuis\s+(\d{1,2})\s*ans?"]:
                    m = re.search(pat, texte, re.IGNORECASE)
                    if m: annees_cv = max(annees_cv, int(m.group(1))); break
                if annees_cv == 0:
                    yrs = re.findall(r"20(\d{2})", texte)
                    if len(yrs) >= 2:
                        annees_cv = max(int(y) for y in yrs) - min(int(y) for y in yrs)

                annees_req = 0
                em = re.search(r"(\d{1,2})\s*/\s*(\d{1,2})\s*ans?", cdc_txt[:4000], re.IGNORECASE)
                if em: annees_req = int(em.group(2))
                elif re.search(r"6.8\s*ans|six.huit", cdc_txt[:4000], re.IGNORECASE): annees_req = 6
                elif re.search(r"senior|expert|confirm", cdc_txt[:4000], re.IGNORECASE): annees_req = 5
                elif re.search(r"junior|debutant", cdc_txt[:4000], re.IGNORECASE): annees_req = 2

                if annees_req == 0: sc_exp_yr = 55
                elif annees_cv == 0: sc_exp_yr = 40
                elif annees_cv >= annees_req: sc_exp_yr = min(80, 60 + (annees_cv - annees_req) * 3)
                else: sc_exp_yr = max(15, int((annees_cv / annees_req) * 60))

                # ── 3. SECTEUR (15%) ─────────────────────────────────────────
                secteurs_ao = crit_in.get("secteurs") or []
                sc_sect = 10  # defaut bas - secteur doit vraiment matcher
                for s_ao in secteurs_ao:
                    s_n = normaliser(s_ao)
                    if any(s_n in normaliser(s) for s in secteurs): sc_sect = 100; break
                    elif s_n in haystack: sc_sect = max(sc_sect, 60)
                # Bonus si entreprise cliente dans le CDC
                for ent in entreprises:
                    if normaliser(ent.get("nom","")) in normaliser(cdc_txt[:2000]):
                        sc_sect = min(100, sc_sect + 20)

                # ── 4. POSTE (10%) ───────────────────────────────────────────
                poste_ao_n = normaliser(crit_in.get("poste") or "")
                mots_p = [m for m in poste_ao_n.split() if len(m) > 3]
                sc_poste = min(100, int(
                    sum(1 for m in mots_p if m in haystack) / max(len(mots_p),1) * 100
                )) if mots_p else 50

                # ── Score global ─────────────────────────────────────────────
                global_score = int(
                    sc_comp    * 0.50 +
                    sc_exp_yr  * 0.25 +
                    sc_sect    * 0.15 +
                    sc_poste   * 0.10
                )
                return {
                    "competences": sc_comp,
                    "experience":  sc_exp_yr,
                    "secteur":     sc_sect,
                    "poste":       sc_poste,
                    "localisation": 75,  # Paris par defaut si pas detecte
                    "global":      global_score,
                }

            # ── SCORER TOUS LES PROFILS puis trier ───────────
            cdc_stock = st.session_state.get("ao_cdc", "")
            criteres_detail = {}
            for s in scores_ao:
                cd = scorer_criteres(s, cdc_stock, crit)
                criteres_detail[s["base"]] = cd
                s["score"] = cd["global"]

            scores_ao.sort(key=lambda x: -x["score"])
            top5 = scores_ao[:5]

            # ── DEBUG : tous les profils avec leur score ──────
            with st.expander("🔍 Voir tous les scores", expanded=False):
                for s in scores_ao:
                    cd   = criteres_detail.get(s["base"], {})
                    info_d = db.get(s["base"], {})
                    texte_d = entrees_cvs.get(s["base"], {}).get("texte","") or ""
                    g = cd.get("global", 0)
                    bg = "#1A6B45" if g>=70 else "#F59E0B" if g>=40 else "#EF4444"
                    st.markdown(
                        f'<span style="background:{bg};color:white;padding:2px 8px;'
                        f'border-radius:6px;font-weight:700">{g}%</span> ' +
                        f'**{s["nom"]}** | '
                        f'Compét:{cd.get("competences",0)} Exp:{cd.get("experience",0)} '
                        f'Profondeur:{cd.get("secteur",0)}',
                        unsafe_allow_html=True
                    )
                    st.caption(
                        f'Compétences DB: {info_d.get("competences",[])} | '
                        f'Texte (200c): {texte_d[:200]}'
                    )

            # ── TABLEAU : SCORE EN PREMIER ────────────────────
            h_cols = st.columns([0.8, 2.2, 1.3, 1.3, 1.3, 1.3, 0.7])
            for col, label in zip(h_cols,
                ["Score", "Candidat", "🎯 Compét.", "📅 Expér.", "🏢 Secteur", "💼 Poste", "CV"]):
                col.markdown(
                    f'<span style="font-size:0.72rem;font-weight:700;color:var(--muted);'
                    f'text-transform:uppercase;letter-spacing:0.4px">{label}</span>',
                    unsafe_allow_html=True)
            st.markdown('<hr style="margin:0.3rem 0 0.5rem">', unsafe_allow_html=True)

            for s in top5:
                cd = criteres_detail[s["base"]]
                row = st.columns([0.8, 2.2, 1.3, 1.3, 1.3, 1.3, 0.7])
                with row[0]:
                    g = cd["global"]
                    if g >= 70:
                        bg, fg = "#1A6B45", "white"
                    elif g >= 40:
                        bg, fg = "#F59E0B", "white"
                    else:
                        bg, fg = "#EF4444", "white"
                    st.markdown(
                        f'<div style="background:{bg};color:{fg};font-size:1.2rem;'
                        f'font-weight:800;padding:6px 10px;border-radius:8px;'
                        f'text-align:center;display:inline-block">{g}%</div>',
                        unsafe_allow_html=True)
                with row[1]:
                    st.markdown(f"**{s['nom']}**")
                    st.caption((s.get("poste","") or "")[:32])
                with row[2]: st.write(nb_etoiles(cd["competences"]))
                with row[3]: st.write(nb_etoiles(cd["experience"]))
                with row[4]: st.write(nb_etoiles(cd["secteur"]))
                with row[5]: st.write(nb_etoiles(cd["poste"]))
                with row[6]:
                    if s.get("pdf"):
                        if st.button("📄", key="ao_cv_" + s["base"], help="Voir CV"):
                            cur = st.session_state.get("ao_cv_ouvert")
                            st.session_state["ao_cv_ouvert"] = None if cur == s["base"] else s["base"]
                            st.rerun()
                if st.session_state.get("ao_cv_ouvert") == s["base"] and s.get("pdf"):
                    afficher_cv_inline(s["pdf"], s["nom"], "ao_" + s["base"], session_key="ao_cv_ouvert")

            # ── DEBUG : tous les scores ───────────────────────
            with st.expander("🔍 Voir tous les scores", expanded=False):
                for s in scores_ao:
                    cd = criteres_detail.get(s["base"], {})
                    info_d = db.get(s["base"], {})
                    texte_d = (entrees_cvs.get(s["base"], {}).get("texte","") or "")[:100]
                    g = s.get("score", 0)
                    c = "#1A6B45" if g >= 70 else "#F59E0B" if g >= 40 else "#6B7280"
                    st.markdown(
                        f'<b style="color:{c}">{g}%</b> — **{s["nom"]}** '
                        f'| Compét: {cd.get("competences",0)} '
                        f'| Exp: {cd.get("experience",0)} '
                        f'| Secteur: {cd.get("secteur",0)}',
                        unsafe_allow_html=True)
                    comps = info_d.get("competences", [])
                    if not comps:
                        st.caption("⚠️ Aucune compétence extraite — relancer IA")
                    elif not texte_d:
                        st.caption("⚠️ Texte PDF vide")
                    else:
                        st.caption("Compétences: " + ", ".join(comps[:5]))

            st.markdown('<hr style="margin:0.5rem 0">', unsafe_allow_html=True)

            # ── SELECTION + EMAILS ────────────────────────────
            noms_top5 = [s["nom"] + " — " + str(s.get("score", 0)) + "%" for s in top5]
            sel_ao = st.multiselect("Sélectionner les profils à contacter :",
                noms_top5, default=noms_top5[:3], key="sel_ao_profils")

            if sel_ao and st.button("✉️ Générer les emails personnalisés", type="primary", key="gen_emails_ao"):
                profils_sel_ao = [s for s in top5
                    if (s["nom"] + " — " + str(s.get("score", 0)) + "%") in sel_ao]
                st.session_state["ao_profils_sel"] = profils_sel_ao
                st.session_state["ao_generer"] = True

    # ========================
    # RESULTATS PLEINE LARGEUR
    # ========================

    # ---- ZONE UPLOAD CV (remontee) ----
    with st.expander("➕ Ajouter des CVs", expanded=False):
        st.caption("Déposez vos fichiers PPTX ou PDF — ils seront copiés dans le dossier Drive et analysés automatiquement par IA")
        uploaded_files = st.file_uploader(
            "Choisir des fichiers",
            type=["pptx", "pdf"],
            accept_multiple_files=True,
            key="upload_cvs"
        )
        if uploaded_files:
            nb_ok = nb_skip = nb_err = 0
            barre_up = st.progress(0)
            for i, uf in enumerate(uploaded_files):
                ext      = os.path.splitext(uf.name)[1].lower()
                dest_dir = PATH_PDF if ext == ".pdf" else PATH_DOSSIER
                dest     = os.path.join(dest_dir, uf.name)
                if os.path.exists(dest):
                    nb_skip += 1
                    continue
                try:
                    data_uf = uf.read()
                    with open(dest, "wb") as fd:
                        fd.write(data_uf)
                    # Si PPTX → convertir en PDF automatiquement
                    if ext == ".pptx":
                        soffice = None
                        for c in ["/Applications/LibreOffice.app/Contents/MacOS/soffice",
                                   shutil.which("soffice"), shutil.which("libreoffice")]:
                            if c and os.path.exists(str(c)):
                                soffice = c; break
                        if soffice:
                            subprocess.run([soffice, "--headless", "--convert-to", "pdf",
                                "--outdir", PATH_PDF, dest],
                                capture_output=True, timeout=30)
                    nb_ok += 1
                except Exception as e:
                    nb_err += 1
                barre_up.progress((i+1)/len(uploaded_files))
            barre_up.empty()
            if nb_ok:
                charger_tous_les_cvs.clear()
                msg = str(nb_ok) + " CV(s) ajouté(s)"
                if nb_err: msg += " · " + str(nb_err) + " erreur(s)"
                st.success(msg + " — Cliquez 'IA' pour analyser")
            if nb_skip:
                st.info(str(nb_skip) + " fichier(s) déjà présent(s) ignoré(s).")

    st.divider()

    # RESULTATS RECHERCHE
    if query:
        resultats = filtrer_cvs(entrees_cvs, query, db)

        if resultats:
            st.write("### " + str(len(resultats)) + " resultat(s)")

            df = pd.DataFrame([{k: v for k, v in r.items() if not k.startswith("_") and k != "Fichier"} for r in resultats])
            # Reordonner colonnes : Score en premier
            cols_ordre = ["Selec.","Score","Nom","Prenom","Telephone","Secteur","Poste","Competences"]
            cols_dispo = [c for c in cols_ordre if c in df.columns]
            df = df[cols_dispo]
            edited_df = st.data_editor(
                df,
                column_config={
                    "Selec.":      st.column_config.CheckboxColumn("✓"),
                    "Score":       st.column_config.TextColumn("Score", width="small"),
                    "Secteur":     st.column_config.TextColumn("Secteur"),
                    "Competences": st.column_config.TextColumn("Competences", width="large"),
                },
                hide_index=True,
                use_container_width=True,
            )

            indices_sel = edited_df[edited_df["Selec."] == True].index.tolist()
            selection   = [resultats[i] for i in indices_sel]

            # Apercu PDF
            if selection:
                for row in selection:
                    pdf_path = row["_pdf_path"]
                    if pdf_path and os.path.exists(pdf_path):
                        st.divider()
                        col_t, col_dl = st.columns([4, 1])
                        with col_t:
                            st.subheader("CV — " + row["Prenom"] + " " + row["Nom"])
                        with col_dl:
                            with open(pdf_path, "rb") as f:
                                st.download_button(
                                    label="Telecharger PDF",
                                    data=f,
                                    file_name=os.path.basename(pdf_path),
                                    mime="application/pdf",
                                    key="dl_" + row["Fichier"],
                                )
                        with open(pdf_path, "rb") as f:
                            b64 = base64.b64encode(f.read()).decode()
                        pdfjs_html = """
<script src="https://cdnjs.cloudflare.com/ajax/libs/pdf.js/3.11.174/pdf.min.js"></script>
<div id="pdf-container" style="background:#f0f0f0;padding:8px;border-radius:8px;overflow-y:auto;max-height:800px;"></div>
<script>
pdfjsLib.GlobalWorkerOptions.workerSrc='https://cdnjs.cloudflare.com/ajax/libs/pdf.js/3.11.174/pdf.worker.min.js';
var b64='""" + b64 + """';
var bin=atob(b64);
var arr=new Uint8Array(bin.length);
for(var i=0;i<bin.length;i++) arr[i]=bin.charCodeAt(i);
pdfjsLib.getDocument({data:arr}).promise.then(function(pdf){
    var container=document.getElementById('pdf-container');
    for(var p=1;p<=pdf.numPages;p++){
        (function(pageNum){
            pdf.getPage(pageNum).then(function(page){
                var scale=1.5;
                var viewport=page.getViewport({scale:scale});
                var canvas=document.createElement('canvas');
                canvas.width=viewport.width;
                canvas.height=viewport.height;
                canvas.style.width='100%';
                canvas.style.marginBottom='8px';
                canvas.style.borderRadius='4px';
                canvas.style.background='white';
                container.appendChild(canvas);
                page.render({canvasContext:canvas.getContext('2d'),viewport:viewport});
            });
        })(p);
    }
});
</script>"""
                        st.components.v1.html(pdfjs_html, height=850, scrolling=True)

                # ---- SECTION ENVOI MULTI-CANAL ----
                st.divider()
                st.subheader("Contacter les candidats")

                msg = st.text_area("Message (utilise pour WA et Email) :", "Bonjour,\n\nJe me permets de vous contacter suite a la reception de votre CV...")
                objet_email = st.text_input("Objet email :", "Opportunite de mission - Noovee")

                canal_col1, canal_col2, canal_col3 = st.columns(3)

                with canal_col1:
                    st.markdown("**📲 WhatsApp**")
                with canal_col2:
                    st.markdown("**✉️ Email (Outlook)**")
                with canal_col3:
                    st.markdown("**🔗 LinkedIn**")

                valides_wa    = [r for r in selection if tel_valide(db.get(r["Fichier"], {}).get("telephone", NO_DATA))]
                valides_email = [r for r in selection if db.get(r["Fichier"], {}).get("email", NO_DATA) not in (NO_DATA, "", "None")]
                tous          = selection

                # Construction HTML multi-canal
                html_canaux = """
<style>
* { font-family: system-ui, sans-serif; box-sizing: border-box; margin:0; padding:0; }
.grid { display:flex; flex-direction:column; gap:8px; padding:4px 0; }
.row { display:grid; grid-template-columns:1fr 1fr 1fr; gap:8px; align-items:center; }
.nom { font-size:13px; font-weight:500; color:#333; padding:6px 0; border-bottom:1px solid #eee; }
.btn { display:flex; align-items:center; gap:6px; padding:8px 12px; border-radius:8px;
       font-size:13px; font-weight:500; text-decoration:none; cursor:pointer; }
.btn-wa   { background:#25D366; color:white; }
.btn-wa:hover { background:#1da851; }
.btn-mail { background:#0078d4; color:white; }
.btn-mail:hover { background:#005a9e; }
.btn-li   { background:#0A66C2; color:white; }
.btn-li:hover { background:#084d93; }
.na { font-size:12px; color:#aaa; padding:8px 0; }
</style>
<div class="grid">
"""
                for row in tous:
                    info = db.get(row["Fichier"], {})
                    tel   = info.get("telephone", NO_DATA)
                    email = info.get("email", NO_DATA)
                    nom_c = (row.get("Prenom","") + " " + row.get("Nom","")).strip() or row["Fichier"][:20]
                    nom_url = urllib.parse.quote(nom_c)

                    # WA
                    if tel_valide(tel):
                        lien_wa = lien_whatsapp(tel, msg)
                        btn_wa = '<a class="btn btn-wa" href="' + lien_wa + '" target="_blank">📲 WA (' + tel + ')</a>'
                    else:
                        btn_wa = '<span class="na">Pas de tel</span>'

                    # Email Outlook
                    if email not in (NO_DATA, "", "None"):
                        corps = urllib.parse.quote(msg)
                        objet_enc = urllib.parse.quote(objet_email)
                        lien_mail = "mailto:" + email + "?subject=" + objet_enc + "&body=" + corps
                        btn_mail = '<a class="btn btn-mail" href="' + lien_mail + '">✉️ ' + email[:20] + '...</a>'
                    else:
                        btn_mail = '<span class="na">Pas d email</span>'

                    # LinkedIn
                    lien_li = "https://www.linkedin.com/search/results/people/?keywords=" + nom_url
                    btn_li = '<a class="btn btn-li" href="' + lien_li + '" target="_blank">🔗 ' + nom_c[:18] + '</a>'

                    html_canaux += (
                        '<div class="nom">' + nom_c + '</div>'
                        '<div class="row">' + btn_wa + btn_mail + btn_li + '</div>'
                    )

                html_canaux += "</div>"
                hauteur_canaux = 80 + len(tous) * 80
                st.components.v1.html(html_canaux, height=hauteur_canaux, scrolling=False)

                # Marquer comme envoye
                st.write("")
                for row in valides_wa:
                    if not db.get(row["Fichier"], {}).get("whatsapp_envoye"):
                        if st.button("Marquer WA envoye — " + row.get("Prenom","") + " " + row.get("Nom",""), key="mark_" + row["Fichier"]):
                            db[row["Fichier"]]["whatsapp_envoye"] = True
                            db[row["Fichier"]]["date_wa"] = datetime.now().strftime("%Y-%m-%d %H:%M")
                            sauvegarder_db(db)
                            st.rerun()

                # Contacts sans telephone
                invalides = [r for r in selection if not tel_valide(db.get(r["Fichier"], {}).get("telephone", NO_DATA))]
                for row in invalides:
                    with st.expander("Ajouter le numero de " + row.get("Prenom","") + " " + row.get("Nom","")):
                        tel_s = st.text_input("Numero :", key="tel_" + row["Fichier"], placeholder="0612345678")
                        if tel_s and st.button("Sauvegarder", key="save_" + row["Fichier"]):
                            db[row["Fichier"]]["telephone"] = tel_s.strip()
                            sauvegarder_db(db)
                            st.rerun()
        else:
            st.info("Aucun resultat.")


# ---------------------------------------------------------------------------
# HELPER : afficher un CV PDF inline (appele depuis plusieurs onglets)
# ---------------------------------------------------------------------------
# afficher_cv_inline definie plus haut


# ---------------------------------------------------------------------------
# MODALE EDITION CONTACT
# ---------------------------------------------------------------------------
if st.session_state.get("edit_contact"):
    base_edit = st.session_state["edit_contact"]
    info_edit = db.get(base_edit, {})
    nom_edit  = (info_edit.get("prenom","") + " " + info_edit.get("nom","")).strip() or base_edit[:30]

    @st.dialog("Modifier le contact — " + nom_edit)
    def dialog_edit():
        info = db.get(base_edit, {})
        def val(champ):
            v = info.get(champ, "")
            return "" if v in (NO_DATA, "None", "null", None) else str(v)

        st.markdown("**Informations personnelles**")
        c1, c2 = st.columns(2)
        with c1:
            e_prenom = st.text_input("Prénom", value=val("prenom"))
            e_email  = st.text_input("Email",  value=val("email"))
            e_secteur = st.selectbox("Secteur principal",
                options=["---"] + sorted(SECTEURS_ACTIFS.keys()),
                index=(["---"] + sorted(SECTEURS_ACTIFS.keys())).index(val("secteur"))
                      if val("secteur") in (["---"] + sorted(SECTEURS_ACTIFS.keys())) else 0
            )
        with c2:
            e_nom   = st.text_input("Nom",       value=val("nom"))
            e_tel   = st.text_input("Téléphone", value=val("telephone"))
            e_poste = st.text_input("Poste",     value=val("poste"))

        # Domaines fonctionnels
        st.markdown("**Domaines fonctionnels**")
        dom_actuels = info.get("domaines", [])
        e_domaines = st.multiselect(
            "Sélectionner les domaines (max 3) :",
            options=DOMAINES_FONCTIONNELS,
            default=[d for d in dom_actuels if d in DOMAINES_FONCTIONNELS],
            max_selections=3,
            key="dom_" + base_edit
        )

        # Competences
        st.markdown("**Compétences clés**")
        e_comp = st.text_area("Compétences (séparées par des virgules)",
            value=", ".join(info.get("competences", [])), height=70)

        # Experiences cles (affichage + edition)
        exps = info.get("experiences_cles", [])
        if exps:
            st.markdown("**🏆 Expériences clés extraites par IA**")
            for i, exp in enumerate(exps):
                with st.container():
                    st.markdown(f"**{exp.get('titre','')}**")
                    st.caption(exp.get("description",""))

        # Entreprises extraites
        entreprises = info.get("entreprises", [])
        if entreprises:
            st.markdown("**🏢 Parcours**")
            badges = " ".join(
                f'<span style="background:#EEF2F5;padding:2px 8px;border-radius:10px;'
                f'font-size:0.8rem;margin:2px">{e.get("nom","")} → {e.get("secteur","")}</span>'
                for e in entreprises
            )
            st.markdown(badges, unsafe_allow_html=True)

        # Texte brut du CV (pour vérification)
        texte_brut_cv = info.get("texte_brut", "")
        if texte_brut_cv:
            with st.expander("📄 Texte brut du CV (extrait par l'app)"):
                st.text_area("", value=texte_brut_cv[:3000], height=200,
                             disabled=True, key="txt_brut_" + base_edit)

        st.divider()
        col_s, col_c = st.columns(2)
        with col_s:
            if st.button("Sauvegarder", type="primary", use_container_width=True):
                db[base_edit].update({
                    "prenom":      e_prenom.strip(),
                    "nom":         e_nom.strip(),
                    "email":       e_email.strip(),
                    "telephone":   e_tel.strip(),
                    "poste":       e_poste.strip(),
                    "secteur":     e_secteur,
                    "domaines":    e_domaines,
                    "competences": [c.strip() for c in e_comp.split(",") if c.strip()],
                })
                sauvegarder_db(db)
                st.session_state["edit_contact"] = None
                st.rerun()
        with col_c:
            if st.button("Annuler", use_container_width=True):
                st.session_state["edit_contact"] = None
                st.rerun()

    dialog_edit()


# ---------------------------------------------------------------------------
# ONGLET 2 — BASE DE CONTACTS
# ---------------------------------------------------------------------------
with tab_contacts:

    # Barre outils
    st.markdown('<div style="display:flex;align-items:center;gap:0.5rem;margin-bottom:0.8rem;">', unsafe_allow_html=True)
    bc1, bc2, bc3, bc4, bc5 = st.columns([1.2, 1.2, 1.5, 1.5, 2])
    with bc1:
        if st.button("Actualiser", key="act_contacts", help="Recharger les CVs"):
            charger_tous_les_cvs.clear()
            st.rerun()
    with bc2:
        if st.button("Convertir PDF", key="conv_contacts", help="Convertir PPTX en PDF"):
            soffice = None
            for c in ["/Applications/LibreOffice.app/Contents/MacOS/soffice",
                       shutil.which("soffice"), shutil.which("libreoffice")]:
                if c and os.path.exists(c):
                    soffice = c
                    break
            if not soffice:
                st.error("LibreOffice non installe.")
            else:
                os.makedirs(PATH_PDF, exist_ok=True)
                fichiers = [f for f in os.listdir(PATH_DOSSIER) if f.endswith(".pptx")]
                nb_ok = nb_err = 0
                prog = st.progress(0)
                for i, f in enumerate(fichiers):
                    cible = os.path.join(PATH_PDF, f.replace(".pptx", ".pdf"))
                    if not os.path.exists(cible):
                        try:
                            r = subprocess.run(
                                [soffice, "--headless", "--convert-to", "pdf",
                                 "--outdir", PATH_PDF, os.path.join(PATH_DOSSIER, f)],
                                capture_output=True, timeout=30
                            )
                            nb_ok += 1 if r.returncode == 0 else 0
                            nb_err += 0 if r.returncode == 0 else 1
                        except:
                            nb_err += 1
                    else:
                        nb_ok += 1
                    prog.progress((i + 1) / len(fichiers))
                prog.empty()
                st.success(str(nb_ok) + " PDF(s) generes")
                charger_tous_les_cvs.clear()
                st.rerun()
    with bc3:
        if st.button("Renommer CVs", key="ren_contacts", help="Format NOM-Prenom-Poste-MMAA"):
            def ascii_clean(s):
                if not s:
                    return ""
                try:
                    return re.sub(r"[^a-zA-Z0-9]", "",
                        unicodedata.normalize("NFD", str(s)).encode("ascii", "ignore").decode())
                except:
                    return ""

            def construire_nom(base, info, noms_utilises):
                nom    = ascii_clean(info.get("nom", "") or "")
                prenom = ascii_clean(info.get("prenom", "") or "")
                poste  = str(info.get("poste", "") or "")
                if poste in (NO_DATA, "null", "None", ""):
                    poste = ""
                mots    = [m for m in re.split(r"[\s\-/,]+", poste) if len(m) > 2][:3]
                poste_c = "".join(ascii_clean(m.capitalize()) for m in mots)
                date_s  = datetime.now().strftime("%m%y")
                parties = [p for p in [nom, prenom, poste_c, date_s] if p]
                if not parties:
                    # Fallback sur le nom de fichier nettoye
                    base_net = ascii_clean(base)[:30]
                    parties  = [base_net, date_s] if base_net else [date_s]
                base_n = "-".join(parties)
                final = base_n
                cnt = 2
                while final in noms_utilises:
                    final = base_n + "_" + str(cnt)
                    cnt += 1
                noms_utilises.add(final)
                return final

            noms_u = set()
            nb_ok = nb_err = 0
            mapping = {}
            previews = []
            for base, info in list(db.items()):
                src = os.path.join(PATH_DOSSIER, base + ".pptx")
                if not os.path.exists(src):
                    continue
                nouveau = construire_nom(base, info, noms_u)
                if nouveau == base:
                    continue
                previews.append(base + " → " + nouveau)
                try:
                    os.rename(src, os.path.join(PATH_DOSSIER, nouveau + ".pptx"))
                    for dp in [PATH_PDF, PATH_DOSSIER]:
                        sp = os.path.join(dp, base + ".pdf")
                        if os.path.exists(sp):
                            os.rename(sp, os.path.join(dp, nouveau + ".pdf"))
                    mapping[base] = nouveau
                    nb_ok += 1
                except Exception as e:
                    nb_err += 1
                    st.warning("Erreur : " + base + " — " + str(e))
            if mapping:
                for a, n in mapping.items():
                    if a in db:
                        db[n] = db.pop(a)
                sauvegarder_db(db)
                charger_tous_les_cvs.clear()
            if nb_ok:
                st.success(str(nb_ok) + " fichier(s) renomme(s)")
                for p in previews[:10]:
                    st.caption(p)
            else:
                st.info("Aucun fichier a renommer (lancez d abord l analyse IA pour remplir les noms).")
            st.rerun()
    with bc4:
        non_enrichis_c = [f for f, v in db.items() if not v.get("ia_enrichi")]
        if st.button(
            "IA (" + str(len(non_enrichis_c)) + ")",
            type="primary",
            key="ia_contacts",
            disabled=(len(non_enrichis_c) == 0 or not AI_KEY),
            help="Analyser les CVs non encore traites"
        ):
            barre = st.progress(0, text="Analyse IA...")
            for i, fname in enumerate(non_enrichis_c):
                # Utiliser le texte stocké en DB en priorité (plus complet et fiable)
                texte = (db[fname].get("texte_brut") or
                         entrees_cvs.get(fname, {}).get("texte", "") or "")
                # Si texte absent → relire le fichier et stocker
                if not texte:
                    data_cv = entrees_cvs.get(fname, {})
                    texte = data_cv.get("texte", "")
                    if texte:
                        db[fname]["texte_brut"] = texte[:8000]
                infos = analyser_cv_avec_ia(texte)
                if infos:
                    for champ in ["nom", "prenom", "email", "poste"]:
                        val = str(infos.get(champ) or "").strip()
                        if val and val.lower() not in ("null", "none", ""):
                            db[fname][champ] = val
                    tel = str(infos.get("telephone") or "").strip()
                    if tel_valide(tel):
                        db[fname]["telephone"] = tel
                    comps = infos.get("competences") or []
                    if isinstance(comps, list) and comps:
                        db[fname]["competences"] = [c for c in comps if c and str(c).lower() != "null"]
                    entreprises_ia = infos.get("entreprises") or []
                    if entreprises_ia:
                        db[fname]["entreprises"] = [
                            {"nom": str(e.get("nom","")).strip(), "secteur": str(e.get("secteur","")).strip()}
                            for e in entreprises_ia
                            if e.get("nom") and str(e.get("nom","")).lower() not in ("null","none","")
                        ]
                secteurs_ia    = list({e["secteur"] for e in db[fname].get("entreprises",[]) if e.get("secteur") and e["secteur"] != "null"})
                secteurs_table = deduire_secteurs_depuis_table(texte)
                tous_secteurs  = list(dict.fromkeys(secteurs_ia + secteurs_table))
                db[fname]["secteurs"] = tous_secteurs[:5]
                db[fname]["secteur"]  = tous_secteurs[0] if tous_secteurs else NO_DATA
                db[fname]["ia_enrichi"] = True
                sauvegarder_db(db)
                barre.progress((i + 1) / len(non_enrichis_c), text=fname[:40])
            barre.empty()
            st.success("Analyse IA terminee !")
            st.rerun()
    with bc5:
        search_contacts = st.text_input("Filtrer...", "", key="search_contacts", placeholder="nom, secteur, poste...")

    st.markdown(
        f'<p style="color:var(--muted);font-size:0.82rem;margin:0.3rem 0 0.8rem;">'
        f'<b style="color:var(--g-dark)">{len(data_list if "data_list" in dir() else db)}</b> contact(s) affiché(s) sur {len(db)} total</p>',
        unsafe_allow_html=True
    )

    if db:
        # Construction de la liste avec PDF
        data_list = []
        for f, v in db.items():
            pdf_dispo = None
            for chemin_pdf in [os.path.join(PATH_PDF, f + ".pdf"), os.path.join(PATH_DOSSIER, f + ".pdf")]:
                if os.path.exists(chemin_pdf):
                    pdf_dispo = chemin_pdf
                    break
            data_list.append({
                "Nom":         v.get("nom", NO_DATA),
                "Prenom":      v.get("prenom", NO_DATA),
                "Email":       v.get("email", NO_DATA),
                "Telephone":   v.get("telephone", NO_DATA),
                "Domaines":    " · ".join(v.get("domaines", [])) or NO_DATA,
                "Secteur":     v.get("secteur", NO_DATA),
                "Poste":       v.get("poste", NO_DATA),
                "Competences": ", ".join(v.get("competences", [])),
                "WA":          "Oui" if v.get("whatsapp_envoye") else "Non",
                "Fichier":     f,
                "_pdf_path":   pdf_dispo,
            })

        # Filtrage
        if search_contacts:
            termes = normaliser(search_contacts).split()
            data_list = [d for d in data_list if all(
                t in normaliser(" ".join(str(val) for val in d.values())) for t in termes
            )]

        st.caption(str(len(data_list)) + " contact(s)")

        # En-tete du tableau
        cols_header = st.columns([2, 1, 1.5, 1.8, 1.5, 1.5, 2.5, 0.6, 0.6])
        headers = ["Nom / Prénom", "CV", "Téléphone", "Email", "Secteur", "Poste", "Compétences", "WA", ""]
        for col, h in zip(cols_header, headers):
            col.markdown(
                f'<span style="font-size:0.75rem;font-weight:600;color:var(--muted);'
                f'text-transform:uppercase;letter-spacing:0.5px">{h}</span>',
                unsafe_allow_html=True
            )
        st.markdown('<hr style="margin:0.4rem 0 0.6rem">', unsafe_allow_html=True)

        # Gestion du CV ouvert
        if "cv_ouvert" not in st.session_state:
            st.session_state["cv_ouvert"] = None

        # Lignes du tableau
        for d in data_list:
            cols = st.columns([2, 1, 1.5, 1.8, 1.5, 1.5, 2.5, 0.6, 0.6])
            
            # Nom affiché - fallback sur le nom de fichier si pas de nom
            nom = d["Nom"] if d["Nom"] != NO_DATA else ""
            prenom = d["Prenom"] if d["Prenom"] != NO_DATA else ""
            if nom or prenom:
                nom_affiche = (prenom + " " + nom).strip()
            else:
                # Extraire depuis le nom de fichier
                parts = d["Fichier"].replace("_", " ").replace("-", " ").split()
                nom_affiche = " ".join(p for p in parts if len(p) > 2)[:30] or d["Fichier"][:20]

            with cols[0]:
                st.write("**" + nom_affiche + "**")

            with cols[1]:
                if d.get("_pdf_path"):
                    if st.button("📄 CV", key="cv_btn_" + d["Fichier"], help="Visualiser le CV"):
                        current = st.session_state.get("cv_ouvert")
                        st.session_state["cv_ouvert"] = None if current == d["Fichier"] else d["Fichier"]
                        st.rerun()
                else:
                    st.caption("—")
            # Affichage CV inline (pleine largeur sous la ligne)
            if st.session_state.get("cv_ouvert") == d["Fichier"] and d.get("_pdf_path"):
                afficher_cv_inline(d["_pdf_path"], nom_affiche, d["Fichier"], session_key="cv_ouvert")

            with cols[2]: st.write(d["Telephone"] if d["Telephone"] != NO_DATA else "")
            with cols[3]: 
                email = d["Email"] if d["Email"] != NO_DATA else ""
                st.write(email[:22] + "..." if len(email) > 22 else email)
            with cols[4]: st.write(d["Secteur"] if d["Secteur"] != NO_DATA else "")
            with cols[5]: 
                poste = d["Poste"] if d["Poste"] != NO_DATA else ""
                st.write(poste[:22] + "..." if len(poste) > 22 else poste)
            with cols[6]: st.write(d["Competences"][:45] + "..." if len(d["Competences"]) > 45 else d["Competences"])
            with cols[7]: st.write("✅" if d["WA"] == "Oui" else "")

            with cols[8]:
                if st.button("🗑️", key="del_" + d["Fichier"], help="Supprimer ce contact"):
                    st.session_state["confirm_del"] = d["Fichier"]
                    st.rerun()

            # Edition du nom inline
            if st.session_state.get("edit_contact") == d["Fichier"]:
                info = db.get(d["Fichier"], {})
                with st.form(key="form_" + d["Fichier"]):
                    col_n1, col_n2, col_n3, col_n4 = st.columns([2, 2, 1, 1])
                    with col_n1:
                        e_prenom = st.text_input("Prenom", value=info.get("prenom", ""))
                    with col_n2:
                        e_nom = st.text_input("Nom", value=info.get("nom", ""))
                    with col_n3:
                        saved = st.form_submit_button("OK", type="primary")
                    with col_n4:
                        cancelled = st.form_submit_button("Annuler")

                if saved:
                    db[d["Fichier"]]["prenom"] = e_prenom.strip()
                    db[d["Fichier"]]["nom"]    = e_nom.strip()
                    sauvegarder_db(db)
                    st.session_state["edit_contact"] = None
                    st.rerun()
                if cancelled:
                    st.session_state["edit_contact"] = None
                    st.rerun()
            else:
                if st.button("✏️", key="edit_" + d["Fichier"], help="Modifier le nom"):
                    st.session_state["edit_contact"] = d["Fichier"]
                    st.rerun()

            # Confirmation suppression
            if st.session_state.get("confirm_del") == d["Fichier"]:
                with st.container():
                    st.warning("Supprimer " + nom_affiche + " ?")
                    c_oui, c_non, c_fich = st.columns([1, 1, 2])
                    with c_oui:
                        if st.button("Oui, DB seulement", key="yes_db_" + d["Fichier"]):
                            del db[d["Fichier"]]
                            sauvegarder_db(db)
                            st.session_state["confirm_del"] = None
                            st.rerun()
                    with c_non:
                        if st.button("Oui + fichiers", key="yes_fich_" + d["Fichier"], type="primary"):
                            del db[d["Fichier"]]
                            for chemin in [
                                os.path.join(PATH_DOSSIER, d["Fichier"] + ".pptx"),
                                os.path.join(PATH_DOSSIER, d["Fichier"] + ".pdf"),
                                os.path.join(PATH_PDF, d["Fichier"] + ".pdf"),
                            ]:
                                if os.path.exists(chemin):
                                    os.remove(chemin)
                            sauvegarder_db(db)
                            charger_tous_les_cvs.clear()
                            st.session_state["confirm_del"] = None
                            st.rerun()
                    with c_fich:
                        if st.button("Annuler", key="no_" + d["Fichier"]):
                            st.session_state["confirm_del"] = None
                            st.rerun()

            # Affichage CV si clique
            if st.session_state.get("cv_ouvert") == d["Fichier"] and d.get("_pdf_path"):
                pdf_path = d["_pdf_path"]
                col_t, col_dl = st.columns([4, 1])
                with col_t:
                    st.subheader("CV — " + nom_affiche)
                with col_dl:
                    with open(pdf_path, "rb") as f_pdf:
                        st.download_button(
                            label="Telecharger",
                            data=f_pdf,
                            file_name=os.path.basename(pdf_path),
                            mime="application/pdf",
                            key="dl_c_" + d["Fichier"],
                        )
                with open(pdf_path, "rb") as f_pdf:
                    b64 = base64.b64encode(f_pdf.read()).decode()
                pdfjs = (
                    "<script src='https://cdnjs.cloudflare.com/ajax/libs/pdf.js/3.11.174/pdf.min.js'></script>"
                    "<div id='pdfc" + d["Fichier"][:8] + "' style='background:#f0f0f0;padding:8px;border-radius:8px;overflow-y:auto;max-height:800px;'></div>"
                    "<script>"
                    "pdfjsLib.GlobalWorkerOptions.workerSrc='https://cdnjs.cloudflare.com/ajax/libs/pdf.js/3.11.174/pdf.worker.min.js';"
                    "var b64='" + b64 + "';"
                    "var bin=atob(b64);var arr=new Uint8Array(bin.length);"
                    "for(var i=0;i<bin.length;i++) arr[i]=bin.charCodeAt(i);"
                    "pdfjsLib.getDocument({data:arr}).promise.then(function(pdf){"
                    "var c=document.getElementById('pdfc" + d["Fichier"][:8] + "');"
                    "for(var p=1;p<=pdf.numPages;p++){(function(pn){pdf.getPage(pn).then(function(page){"
                    "var vp=page.getViewport({scale:1.5});var cv=document.createElement('canvas');"
                    "cv.width=vp.width;cv.height=vp.height;cv.style.width='100%';"
                    "cv.style.marginBottom='8px';cv.style.background='white';"
                    "c.appendChild(cv);page.render({canvasContext:cv.getContext('2d'),viewport:vp});"
                    "});})(p);}});</script>"
                )
                st.components.v1.html(pdfjs, height=850, scrolling=True)
                st.divider()

        # ---- ENVOI MULTI-CANAL depuis la base ----
        st.divider()
        st.subheader("Contacter depuis la base")

        contacts_disponibles = [
            d["Prenom"] + " " + d["Nom"] + " | " + d["Fichier"]
            for d in data_list
        ]
        selection_base = st.multiselect(
            "Selectionner des contacts a contacter :",
            options=contacts_disponibles,
            key="multiselect_contacts"
        )

        if selection_base:
            msg_base   = st.text_area("Message :", "Bonjour,\n\nJe me permets de vous contacter concernant une opportunite de mission.", key="msg_base")
            objet_base = st.text_input("Objet email :", "Opportunite de mission - Noovee", key="objet_base")

            contacts_sel = [d for d in data_list if (d["Prenom"] + " " + d["Nom"] + " | " + d["Fichier"]) in selection_base]

            html_envoi = """
<style>
* { font-family: system-ui, sans-serif; box-sizing: border-box; }
.contact-row { border-bottom: 1px solid #EEF2F5; padding: 10px 0; }
.contact-name { font-weight: 600; color: #1E2A3A; margin-bottom: 8px; font-size: 14px; }
.btns { display: flex; gap: 8px; flex-wrap: wrap; }
.btn { display:inline-flex; align-items:center; gap:6px; padding:7px 14px;
       border-radius:8px; font-size:13px; font-weight:500; text-decoration:none; }
.wa   { background:#25D366; color:white; }
.mail { background:#093F28; color:white; }
.li   { background:#0A66C2; color:white; }
.na   { color:#8898AA; font-size:12px; font-style:italic; }
</style>
"""
            import urllib.parse as up
            for d in contacts_sel:
                info   = db.get(d["Fichier"], {})
                tel    = info.get("telephone", NO_DATA)
                email  = info.get("email", NO_DATA)
                nom_c  = (d["Prenom"] + " " + d["Nom"]).strip()
                nom_url = up.quote(nom_c)

                btn_wa   = ('<a class="btn wa" href="' + lien_whatsapp(tel, msg_base) + '" target="_blank">📲 WhatsApp</a>'
                            if tel_valide(tel) else '<span class="na">Pas de telephone</span>')
                if email not in (NO_DATA, "", "None"):
                    lien_m = "mailto:" + email + "?subject=" + up.quote(objet_base) + "&body=" + up.quote(msg_base)
                    btn_mail = '<a class="btn mail" href="' + lien_m + '">✉️ Outlook</a>'
                else:
                    btn_mail = '<span class="na">Pas d email</span>'
                btn_li = '<a class="btn li" href="https://www.linkedin.com/search/results/people/?keywords=' + nom_url + '" target="_blank">🔗 LinkedIn</a>'

                html_envoi += (
                    '<div class="contact-row">'
                    '<div class="contact-name">' + nom_c + '</div>'
                    '<div class="btns">' + btn_wa + btn_mail + btn_li + '</div>'
                    '</div>'
                )

            st.components.v1.html(html_envoi, height=80 + len(contacts_sel) * 80, scrolling=False)

        st.divider()
        st.subheader("Ajouter un contact manuellement")
        c1, c2, c3, c4 = st.columns(4)
        with c1: new_nom    = st.text_input("Nom", placeholder="DUPONT")
        with c2: new_prenom = st.text_input("Prenom", placeholder="Marie")
        with c3: new_email  = st.text_input("Email", placeholder="marie@email.com")
        with c4: new_tel    = st.text_input("Telephone", placeholder="0612345678")

        if st.button("Ajouter le contact"):
            if new_nom and new_prenom:
                key = "MANUEL-" + new_nom.upper() + " " + new_prenom
                db[key] = {
                    "nom": new_nom.upper(), "prenom": new_prenom,
                    "email": new_email.strip(), "telephone": new_tel.strip(),
                    "poste": NO_DATA, "secteur": NO_DATA, "competences": [],
                    "ia_enrichi": False,
                    "date_ajout": datetime.now().strftime("%Y-%m-%d %H:%M"),
                    "whatsapp_envoye": False,
                }
                sauvegarder_db(db)
                st.success("Contact ajoute !")
                st.rerun()
            else:
                st.warning("Nom et Prenom obligatoires.")


    # Affichage CV AO (pleine largeur)
    if st.session_state.get("ao_cv_ouvert"):
        base_cv_ao = st.session_state["ao_cv_ouvert"]
        pdf_cv_ao = None
        for cp in [os.path.join(PATH_PDF, base_cv_ao + ".pdf"),
                    os.path.join(PATH_DOSSIER, base_cv_ao + ".pdf")]:
            if os.path.exists(cp):
                pdf_cv_ao = cp
                break
        if pdf_cv_ao:
            info_ao = db.get(base_cv_ao, {})
            nom_ao  = (info_ao.get("prenom","") + " " + info_ao.get("nom","")).strip() or base_cv_ao[:25]
            col_cv1, col_cv2 = st.columns([4, 1])
            with col_cv1:
                st.subheader("CV — " + nom_ao)
            with col_cv2:
                with open(pdf_cv_ao, "rb") as f_ao:
                    st.download_button("Telecharger", data=f_ao,
                        file_name=os.path.basename(pdf_cv_ao),
                        mime="application/pdf", key="dl_ao_" + base_cv_ao)
            with open(pdf_cv_ao, "rb") as f_ao:
                b64_ao = base64.b64encode(f_ao.read()).decode()
            pdfjs_ao = (
                "<script src='https://cdnjs.cloudflare.com/ajax/libs/pdf.js/3.11.174/pdf.min.js'></script>"
                "<div id='pdfao' style='background:#f0f0f0;padding:8px;border-radius:8px;overflow-y:auto;max-height:750px;'></div>"
                "<script>"
                "pdfjsLib.GlobalWorkerOptions.workerSrc='https://cdnjs.cloudflare.com/ajax/libs/pdf.js/3.11.174/pdf.worker.min.js';"
                "var b='"+b64_ao+"';var bn=atob(b);var ar=new Uint8Array(bn.length);"
                "for(var i=0;i<bn.length;i++) ar[i]=bn.charCodeAt(i);"
                "pdfjsLib.getDocument({data:ar}).promise.then(function(pdf){"
                "var c=document.getElementById('pdfao');"
                "for(var p=1;p<=pdf.numPages;p++){(function(pn){pdf.getPage(pn).then(function(page){"
                "var vp=page.getViewport({scale:1.5});var cv=document.createElement('canvas');"
                "cv.width=vp.width;cv.height=vp.height;cv.style.width='100%';"
                "cv.style.marginBottom='8px';cv.style.background='white';"
                "c.appendChild(cv);page.render({canvasContext:cv.getContext('2d'),viewport:vp});"
                "});})(p);}});</script>"
            )
            st.components.v1.html(pdfjs_ao, height=780, scrolling=True)
            st.divider()

    # Emails generes par l AO (pleine largeur)
    if st.session_state.get("ao_generer") and st.session_state.get("ao_criteres"):
        crit   = st.session_state["ao_criteres"]
        client = st.session_state.get("ao_client", "le client")
        comp_ao    = crit.get("competences") or []
        secteur_ao = crit.get("secteurs") or []
        mots_ao    = [normaliser(c) for c in comp_ao + secteur_ao]

        scores_ao2 = []
        for base, data in entrees_cvs.items():
            info    = db.get(base, {})
            haystack = normaliser(data.get("texte","") + " " + " ".join(info.get("competences",[])) + " " + info.get("secteur",""))
            score   = sum(1 for m in mots_ao if m and m in haystack)
            score_p = min(100, int(score / max(len(mots_ao),1) * 100))
            if score > 0:
                nom_s = (info.get("prenom","") + " " + info.get("nom","")).strip() or base[:20]
                scores_ao2.append({
                    "pct": score_p, "nom": nom_s, "base": base,
                    "poste": info.get("poste",""),
                    "competences": ", ".join(info.get("competences",[])[:4]),
                    "tel": info.get("telephone", NO_DATA),
                    "email": info.get("email", NO_DATA),
                })
        scores_ao2.sort(key=lambda x: -x["pct"])

        st.subheader("Emails personnalises — " + (crit.get("poste") or "Mission") + " · " + client)
        for profil in scores_ao2[:3]:
            with st.spinner("Generation pour " + profil["nom"] + "..."):
                prompt_mail = "\n".join([
                    "Tu es recruteur chez Noovee, cabinet de sourcing de consultants experts.",
                    "Redige un email professionnel et personnalise en francais pour " + profil["nom"] + ".",
                    "Son poste : " + profil["poste"] + ". Ses competences : " + profil["competences"] + ".",
                    "Mission : " + (crit.get("poste") or "consultant expert") + " chez " + client + ".",
                    "Duree : " + (crit.get("duree") or "a definir") + ". Lieu : " + (crit.get("localisation") or "Ile de France") + ".",
                    "L email doit : presenter brievement Noovee, expliquer pourquoi ce profil correspond,",
                    "donner 2 details cles sur la mission, inviter a un echange.",
                    "Ton professionnel et chaleureux. 120 mots max.",
                    "Reponds UNIQUEMENT avec le corps de l email."
                ])
                try:
                    prov_e, key_e = get_provider_email()
                    # Claude Sonnet pour meilleure qualite redactionnelle
                    model_e = "claude-sonnet-4-6" if prov_e == "ANTHROPIC" else None
                    email_body = appeler_ia(prompt_mail, prov_e, key_e,
                                           max_tokens=400, model_override=model_e)
                except Exception:
                    email_body = "Erreur generation email."

            objet_ao = "Opportunite de mission - " + (crit.get("poste") or "Expert") + " - Noovee"
            with st.expander(profil["nom"] + " — " + str(profil["pct"]) + "% de matching", expanded=True):
                objet_ed = st.text_input("Objet :", value=objet_ao, key="obj2_" + profil["base"])
                corps_ed = st.text_area("Email :", value=email_body, height=180, key="corps2_" + profil["base"])
                c1, c2, c3 = st.columns(3)
                with c1:
                    if profil["email"] not in (NO_DATA,"","None"):
                        lm = "mailto:" + profil["email"] + "?subject=" + urllib.parse.quote(objet_ed) + "&body=" + urllib.parse.quote(corps_ed)
                        st.link_button("✉️ Outlook", lm)
                    else:
                        st.caption("Pas d email")
                with c2:
                    if tel_valide(profil["tel"]):
                        st.link_button("📲 WhatsApp", lien_whatsapp(profil["tel"], corps_ed))
                with c3:
                    st.link_button("🔗 LinkedIn", "https://www.linkedin.com/search/results/people/?keywords=" + urllib.parse.quote(profil["nom"]))

        st.session_state["ao_generer"] = False

# ---------------------------------------------------------------------------
# ONGLET 3 — DOUBLONS
# ---------------------------------------------------------------------------
with tab_doublons:
    st.subheader("Analyse des doublons")

    def trouver_doublons(db):
        """Detecte les doublons par telephone, email ou nom+prenom similaires."""
        groupes = []
        vus = set()
        items = list(db.items())
        for i, (base1, v1) in enumerate(items):
            if base1 in vus:
                continue
            groupe = [base1]
            tel1   = str(v1.get("telephone", "")).strip()
            email1 = str(v1.get("email", "")).strip().lower()
            nom1   = normaliser(v1.get("nom", "") + " " + v1.get("prenom", ""))
            for base2, v2 in items[i+1:]:
                if base2 in vus:
                    continue
                tel2   = str(v2.get("telephone", "")).strip()
                email2 = str(v2.get("email", "")).strip().lower()
                nom2   = normaliser(v2.get("nom", "") + " " + v2.get("prenom", ""))
                doublon = False
                if tel1 and tel1 != NO_DATA and tel1 == tel2:
                    doublon = True
                if email1 and email1 != NO_DATA and "@" in email1 and email1 == email2:
                    doublon = True
                if nom1 and nom1 not in ("--- ---", "candidat candidat") and nom1 == nom2 and len(nom1) > 5:
                    doublon = True
                # Comparaison par poste + competences identiques (cas CV sans nom)
                poste1 = normaliser(v1.get("poste", ""))
                poste2 = normaliser(v2.get("poste", ""))
                comp1  = normaliser(" ".join(v1.get("competences", [])))
                comp2  = normaliser(" ".join(v2.get("competences", [])))
                if (poste1 and poste1 != normaliser(NO_DATA) and poste1 == poste2
                        and comp1 and comp1 == comp2 and len(comp1) > 10):
                    doublon = True
                if doublon:
                    groupe.append(base2)
                    vus.add(base2)
            if len(groupe) > 1:
                vus.add(base1)
                groupes.append(groupe)
        return groupes

    doublons = trouver_doublons(db)

    if not doublons:
        st.success("Aucun doublon detecte dans la base !")
    else:
        st.warning(str(len(doublons)) + " groupe(s) de doublons detectes")

        for idx, groupe in enumerate(doublons):
            st.divider()
            st.write("**Groupe " + str(idx+1) + "** — " + str(len(groupe)) + " entrees similaires")

            rows_groupe = []
            for base in groupe:
                v = db.get(base, {})
                rows_groupe.append({
                    "Fichier":   base,
                    "Nom":       v.get("nom", NO_DATA),
                    "Prenom":    v.get("prenom", NO_DATA),
                    "Email":     v.get("email", NO_DATA),
                    "Telephone": v.get("telephone", NO_DATA),
                    "Secteur":   v.get("secteur", NO_DATA),
                    "Date":      v.get("date_ajout", ""),
                })
            df_groupe = pd.DataFrame(rows_groupe)
            st.dataframe(df_groupe, hide_index=True, use_container_width=True)

            suppr_fichiers = st.checkbox(
                "Supprimer aussi les fichiers PPTX et PDF du dossier",
                key="del_files_" + str(idx)
            )

            def supprimer_contacts(a_supprimer, avec_fichiers):
                nb_f = 0
                for base in a_supprimer:
                    if base in db:
                        del db[base]
                    if avec_fichiers:
                        for chemin in [
                            os.path.join(PATH_DOSSIER, base + ".pptx"),
                            os.path.join(PATH_DOSSIER, base + ".pdf"),
                            os.path.join(PATH_PDF, base + ".pdf"),
                        ]:
                            if os.path.exists(chemin):
                                os.remove(chemin)
                                nb_f += 1
                return nb_f

            col_g1, col_g2 = st.columns(2)
            with col_g1:
                if st.button("Garder le 1er, supprimer les autres", key="dup_keep_" + str(idx)):
                    a_supprimer = groupe[1:]
                    nb_f = supprimer_contacts(a_supprimer, suppr_fichiers)
                    sauvegarder_db(db)
                    if suppr_fichiers:
                        charger_tous_les_cvs.clear()
                    msg = str(len(a_supprimer)) + " doublon(s) supprimes de la DB"
                    if nb_f:
                        msg += " + " + str(nb_f) + " fichier(s) effaces"
                    st.success(msg + " !")
                    st.rerun()
            with col_g2:
                if st.button("Garder le plus recent", key="dup_recent_" + str(idx)):
                    groupe_tri = sorted(groupe, key=lambda b: db.get(b, {}).get("date_ajout", ""), reverse=True)
                    a_supprimer = groupe_tri[1:]
                    nb_f = supprimer_contacts(a_supprimer, suppr_fichiers)
                    sauvegarder_db(db)
                    if suppr_fichiers:
                        charger_tous_les_cvs.clear()
                    msg = "Contact le plus recent conserve, " + str(len(a_supprimer)) + " doublon(s) supprimes"
                    if nb_f:
                        msg += " + " + str(nb_f) + " fichier(s) effaces"
                    st.success(msg + " !")
                    st.rerun()


# ---------------------------------------------------------------------------
# ONGLET 4 — TABLE SECTEURS (editable)
# ---------------------------------------------------------------------------
with tab_secteurs:
    st.subheader("Table de correspondance Entreprises / Secteurs")
    st.caption("Modifiez, ajoutez ou supprimez des mots-cles pour chaque secteur. Cliquez Sauvegarder pour appliquer.")

    PATH_SECTEURS = os.path.join(os.path.dirname(os.path.abspath(__file__)), "secteurs_custom.json")

    def charger_secteurs_custom():
        """Charge les secteurs depuis le fichier custom si existe, sinon retourne SECTEURS par defaut."""
        if os.path.exists(PATH_SECTEURS):
            try:
                with open(PATH_SECTEURS, "r", encoding="utf-8") as f:
                    return json.load(f)
            except:
                pass
        return {k: list(v) for k, v in SECTEURS.items()}

    def sauvegarder_secteurs_custom(secteurs):
        with open(PATH_SECTEURS, "w", encoding="utf-8") as f:
            json.dump(secteurs, f, ensure_ascii=False, indent=2)

    secteurs_edit = dict(SECTEURS_ACTIFS)  # Deja charge au demarrage

    # Selecteur de secteur a editer
    secteur_choisi = st.selectbox("Choisir un secteur a modifier :", sorted(secteurs_edit.keys()))

    if secteur_choisi:
        mots_actuels = secteurs_edit[secteur_choisi]
        st.write("**" + str(len(mots_actuels)) + " mots-cles** pour ce secteur :")

        # Affichage editable des mots-cles
        mots_texte = st.text_area(
            "Mots-cles (un par ligne) :",
            value="\n".join(sorted(mots_actuels)),
            height=300,
            key="mots_" + secteur_choisi
        )

        col_s1, col_s2 = st.columns(2)
        with col_s1:
            if st.button("Sauvegarder ce secteur", type="primary"):
                nouveaux_mots = [m.strip().lower() for m in mots_texte.strip().splitlines() if m.strip()]
                secteurs_edit[secteur_choisi] = nouveaux_mots
                sauvegarder_secteurs_custom(secteurs_edit)
                # Mettre a jour TABLE_RAPIDE en memoire
                TABLE_RAPIDE.clear()
                for s, mots in secteurs_edit.items():
                    for mot in mots:
                        TABLE_RAPIDE[normaliser(mot)] = s
                st.success("Secteur mis a jour avec " + str(len(nouveaux_mots)) + " mots-cles ! Relancez l app pour appliquer a tous les CVs.")

        with col_s2:
            if st.button("Recalculer tous les secteurs"):
                # Recharge la table avec les nouveaux mots-cles
                nouvelle_table = {}
                for s, mots in secteurs_edit.items():
                    for mot in mots:
                        nouvelle_table[normaliser(mot)] = s
                for base, data in entrees_cvs.items():
                    texte = data.get("texte", "")
                    t = normaliser(texte)
                    scores = {}
                    for mot_norme, secteur in nouvelle_table.items():
                        if mot_norme in t:
                            scores[secteur] = scores.get(secteur, 0) + 1
                    db[base]["secteur"] = max(scores, key=scores.get) if scores else NO_DATA
                sauvegarder_db(db)
                st.success("Secteurs recalcules pour tous les CVs !")
                st.rerun()

    st.divider()
    # Ajouter un nouveau secteur
    st.subheader("Ajouter un nouveau secteur")
    col_n1, col_n2 = st.columns([2, 4])
    with col_n1:
        nouveau_secteur = st.text_input("Nom du secteur", placeholder="Ex: Juridique")
    with col_n2:
        nouveaux_mots_input = st.text_area("Mots-cles (un par ligne)", placeholder="cabinet avocats\njuridique\ndroit\nnotaire", height=120)

    if st.button("Creer ce secteur"):
        if nouveau_secteur and nouveaux_mots_input:
            mots = [m.strip().lower() for m in nouveaux_mots_input.splitlines() if m.strip()]
            secteurs_edit[nouveau_secteur] = mots
            sauvegarder_secteurs_custom(secteurs_edit)
            st.success("Secteur '" + nouveau_secteur + "' cree avec " + str(len(mots)) + " mots-cles !")
            st.rerun()
        else:
            st.warning("Remplissez le nom et les mots-cles.")

    st.divider()
    # Vue complete de la table
    with st.expander("Voir toute la table de correspondance"):
        all_rows = []
        for secteur, mots in sorted(secteurs_edit.items()):
            for mot in sorted(mots):
                all_rows.append({"Secteur": secteur, "Mot-cle / Entreprise": mot})
        if all_rows:
            st.dataframe(pd.DataFrame(all_rows), hide_index=True, use_container_width=True)


# ---------------------------------------------------------------------------
# ONGLET 5 — OUTILS
# ---------------------------------------------------------------------------
with tab_outils:
    st.subheader("Maintenance")

    c1, c2, c3, c4 = st.columns(4)
    c1.metric("CVs indexes", len(entrees_cvs))
    c2.metric("Contacts en DB", len(db))
    c3.metric("Analyses IA", sum(1 for v in db.values() if v.get("ia_enrichi")))
    c4.metric("WA envoyes", sum(1 for v in db.values() if v.get("whatsapp_envoye")))

    if ALL_KEYS:
        labels = {
            "MISTRAL_API_KEY":   "🇫🇷 Mistral (extraction)",
            "ANTHROPIC_API_KEY": "🤖 Claude (emails)",
            "OPENAI_API_KEY":    "⚡ GPT-4o (scoring)",
            "GROQ_API_KEY":      "Groq",
            "GEMINI_API_KEY":    "Gemini",
        }
        for var, key in ALL_KEYS.items():
            st.success(labels.get(var, var) + " — connecté")
    else:
        st.warning("Aucune IA configuree. Ajoutez MISTRAL_API_KEY, OPENAI_API_KEY ou ANTHROPIC_API_KEY dans .env")

    st.divider()

    # Stats par secteur
    st.subheader("Repartition par secteur")
    secteurs_stats = {}
    for v in db.values():
        s = v.get("secteur", NO_DATA)
        secteurs_stats[s] = secteurs_stats.get(s, 0) + 1
    df_secteurs = pd.DataFrame(
        [(s, n) for s, n in sorted(secteurs_stats.items(), key=lambda x: -x[1])],
        columns=["Secteur", "Nombre de CVs"]
    )
    st.dataframe(df_secteurs, hide_index=True, use_container_width=True)

    st.divider()
    col1, col2, col3 = st.columns(3)

    with col1:
        if st.button("Actualiser les CVs"):
            charger_tous_les_cvs.clear()
            st.rerun()

    with col2:
        if st.button("Convertir en PDF"):
            soffice = None
            for c in ["/Applications/LibreOffice.app/Contents/MacOS/soffice",
                       shutil.which("soffice"), shutil.which("libreoffice")]:
                if c and os.path.exists(c):
                    soffice = c
                    break
            if not soffice:
                st.error("LibreOffice non installe.")
            else:
                os.makedirs(PATH_PDF, exist_ok=True)
                fichiers = [f for f in os.listdir(PATH_DOSSIER) if f.endswith(".pptx")]
                nb_ok = nb_err = 0
                prog = st.progress(0)
                for i, f in enumerate(fichiers):
                    cible = os.path.join(PATH_PDF, f.replace(".pptx", ".pdf"))
                    if not os.path.exists(cible):
                        try:
                            r = subprocess.run(
                                [soffice, "--headless", "--convert-to", "pdf",
                                 "--outdir", PATH_PDF, os.path.join(PATH_DOSSIER, f)],
                                capture_output=True, timeout=30
                            )
                            nb_ok += 1 if r.returncode == 0 else 0
                            nb_err += 0 if r.returncode == 0 else 1
                        except:
                            nb_err += 1
                    else:
                        nb_ok += 1
                    prog.progress((i + 1) / len(fichiers))
                st.success(str(nb_ok) + " PDF(s) generes")
                charger_tous_les_cvs.clear()
                st.rerun()

    with col3:
        if st.button("Re-analyser tous les CVs"):
            for base in db:
                db[base]["ia_enrichi"] = False
            sauvegarder_db(db)
            st.success("Pret pour re-analyse.")
            st.rerun()

    st.divider()

    # ---- RENOMMAGE DES CVS ----
    st.subheader("Renommer les CVs")
    st.caption("Format : NOM-Prenom-PosteCle-MMAA.pptx")

    def construire_nouveau_nom(base, info, noms_utilises):
        """Construit le nouveau nom de fichier normalise."""
        nom    = re.sub(r"[^a-zA-Z]", "", unicodedata.normalize("NFD", info.get("nom", "").upper()).encode("ascii","ignore").decode())
        prenom = re.sub(r"[^a-zA-Z]", "", unicodedata.normalize("NFD", info.get("prenom", "").capitalize()).encode("ascii","ignore").decode())
        poste  = info.get("poste", "")
        # Nettoyer le poste : garder les mots significatifs, max 3 mots, CamelCase
        mots_poste = [m for m in re.split("[\\s\\-/,]+", poste) if len(m) > 2][:3]
        poste_clean = "".join([
            re.sub(r"[^a-zA-Z]", "", unicodedata.normalize("NFD", m.capitalize()).encode("ascii","ignore").decode())
            for m in mots_poste
        ])
        date_str = datetime.now().strftime("%m%y")
        # Construire le nom de base
        parties = [p for p in [nom, prenom, poste_clean, date_str] if p]
        nouveau = "-".join(parties)
        if not nouveau:
            nouveau = re.sub(r"[^a-zA-Z0-9\-]", "", base)[:40]
        # Gerer les doublons de noms
        nom_final = nouveau
        compteur = 2
        while nom_final in noms_utilises:
            nom_final = nouveau + "_" + str(compteur)
            compteur += 1
        noms_utilises.add(nom_final)
        return nom_final

    # Preview des renommages
    if st.button("Previsualiser les renommages"):
        noms_utilises = set()
        previews = []
        for base, info in db.items():
            if os.path.exists(os.path.join(PATH_DOSSIER, base + ".pptx")):
                nouveau = construire_nouveau_nom(base, info, noms_utilises)
                if nouveau != base:
                    previews.append({"Avant": base, "Apres": nouveau + ".pptx"})
        if previews:
            st.dataframe(pd.DataFrame(previews), hide_index=True, use_container_width=True)
            st.session_state["previews_renommage"] = previews
        else:
            st.info("Tous les fichiers sont deja au bon format.")

    if st.button("Renommer tous les CVs", type="primary"):
        noms_utilises = set()
        nb_ok = nb_err = 0
        mapping = {}  # ancien_base -> nouveau_base
        for base, info in list(db.items()):
            pptx_src = os.path.join(PATH_DOSSIER, base + ".pptx")
            if not os.path.exists(pptx_src):
                continue
            nouveau_base = construire_nouveau_nom(base, info, noms_utilises)
            if nouveau_base == base:
                continue
            try:
                # Renommer PPTX
                os.rename(pptx_src, os.path.join(PATH_DOSSIER, nouveau_base + ".pptx"))
                # Renommer PDF si existe
                for dossier_pdf in [PATH_PDF, PATH_DOSSIER]:
                    src_pdf = os.path.join(dossier_pdf, base + ".pdf")
                    if os.path.exists(src_pdf):
                        os.rename(src_pdf, os.path.join(dossier_pdf, nouveau_base + ".pdf"))
                mapping[base] = nouveau_base
                nb_ok += 1
            except Exception as e:
                nb_err += 1

        # Mise a jour de la DB avec les nouveaux noms
        if mapping:
            for ancien, nouveau in mapping.items():
                if ancien in db:
                    db[nouveau] = db.pop(ancien)
            sauvegarder_db(db)
            charger_tous_les_cvs.clear()

        st.success(str(nb_ok) + " fichier(s) renomme(s) - " + str(nb_err) + " erreur(s)")
        st.rerun()

    st.divider()
    st.caption("Dossier : " + PATH_DOSSIER)
    st.caption("DB : " + PATH_DB)
